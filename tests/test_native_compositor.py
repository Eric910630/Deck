"""
测试 Native Shape Compiler
验证 DOM 样式提取和原生 PPT 绘制功能
"""

import asyncio
import sys
from pathlib import Path

# 添加项目根目录到路径
sys.path.insert(0, str(Path(__file__).parent.parent))

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Cm
from loguru import logger

from langchain_deck.rendering.dom_analyzer import DOMAnalyzer
from langchain_deck.rendering.native_compositor import NativeCompositor


async def test_native_compositor():
    """测试原生合成器"""
    logger.info("="*80)
    logger.info("测试 Native Shape Compiler")
    logger.info("="*80)
    
    # 1. 准备测试 HTML（使用之前生成的 CSS-First HTML）
    html_file = Path("html_output/test_single_slide_css_first.html")
    if not html_file.exists():
        logger.error(f"测试 HTML 文件不存在: {html_file}")
        logger.info("请先运行 test_single_slide_layout.py 生成测试 HTML")
        return
    
    html_content = html_file.read_text(encoding='utf-8')
    
    # 2. 使用 Playwright 加载 HTML 并提取样式
    logger.info("--- [测试1] 使用 DOMAnalyzer 提取样式...")
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(viewport={'width': 1920, 'height': 1080})
        
        await page.set_content(html_content)
        
        # 提取布局数据
        analyzer = DOMAnalyzer()
        layout_data = await analyzer.extract_layout_data(page)
        
        logger.info(f"--- [测试1] ✅ 提取了 {len(layout_data)} 个元素的样式数据")
        
        # 打印第一个元素的数据结构（用于调试）
        if layout_data:
            logger.info(f"--- [测试1] 示例元素数据:")
            logger.info(f"   ID: {layout_data[0].get('id')}")
            logger.info(f"   Type: {layout_data[0].get('type')}")
            logger.info(f"   Geometry: {layout_data[0].get('geometry')}")
            logger.info(f"   Style keys: {list(layout_data[0].get('style', {}).keys())}")
        
        await browser.close()
    
    # 3. 使用 NativeCompositor 绘制到 PPT
    logger.info("--- [测试2] 使用 NativeCompositor 绘制到 PPT...")
    
    prs = Presentation()
    prs.slide_width = Cm(33.867)  # 16:9
    prs.slide_height = Cm(19.05)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    compositor = NativeCompositor()
    compositor.composite_slide(slide, layout_data)
    
    # 4. 保存 PPT
    output_file = Path("outputs/ppt/test_native_compositor.pptx")
    output_file.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_file))
    
    logger.info(f"--- [测试2] ✅ PPT 已保存到: {output_file}")
    logger.info("="*80)
    logger.info("✅ Native Shape Compiler 测试完成！")
    logger.info("="*80)


if __name__ == "__main__":
    asyncio.run(test_native_compositor())

