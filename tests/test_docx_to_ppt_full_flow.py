#!/usr/bin/env python3
"""
完整流程测试：从Demo文档.docx生成PPT
使用Demo文档.docx的内容作为PPT的主要内容，运行完整流程
带详细探针和日志分析
"""

import asyncio
import json
import sys
from pathlib import Path
from typing import Dict, Any
from datetime import datetime
from loguru import logger
from docx import Document

# 添加项目根目录到路径
sys.path.insert(0, str(Path(__file__).parent.parent))

from ppt_filler import PPTFiller

# 配置日志系统
def setup_logging(log_file: str):
    """配置日志系统，同时输出到控制台和文件"""
    logger.remove()  # 移除默认处理器
    
    # 控制台输出（带颜色）
    logger.add(
        sys.stdout,
        format="<green>{time:YYYY-MM-DD HH:mm:ss}</green> | <level>{level: <8}</level> | <level>{message}</level>",
        level="INFO"
    )
    
    # 文件输出（详细格式）
    logger.add(
        log_file,
        format="{time:YYYY-MM-DD HH:mm:ss} | {level: <8} | {name}:{function}:{line} | {message}",
        level="DEBUG",
        rotation="100 MB",
        retention="10 days"
    )
    
    return logger

def extract_docx_content(docx_path: str) -> Dict[str, Any]:
    """从docx文件中提取文本内容（详细探针）"""
    doc = Document(docx_path)
    content_parts = []
    detailed_info = {
        "paragraphs": [],
        "tables": [],
        "structure": []
    }
    
    print(f"\n{'='*80}")
    print("【详细探针】docx内容提取过程")
    print("="*80)
    
    # 提取段落
    print(f"\n📝 段落提取（共{len(doc.paragraphs)}个段落）:")
    for idx, para in enumerate(doc.paragraphs):
        para_text = para.text.strip()
        if para_text:
            style_name = para.style.name if para.style else "无样式"
            is_bold = any(run.bold for run in para.runs if run.bold)
            font_size = None
            for run in para.runs:
                if run.font.size:
                    font_size = run.font.size.pt
                    break
            
            para_info = {
                "index": idx,
                "text": para_text,
                "style": style_name,
                "is_bold": is_bold,
                "font_size": font_size,
                "length": len(para_text)
            }
            detailed_info["paragraphs"].append(para_info)
            
            print(f"   段落{idx}:")
            print(f"     文本: {para_text[:100]}{'...' if len(para_text) > 100 else ''}")
            print(f"     样式: {style_name}")
            print(f"     加粗: {is_bold}")
            print(f"     字号: {font_size}pt" if font_size else "     字号: 默认")
            print(f"     长度: {len(para_text)}字符")
            
            content_parts.append(para_text)
    
    # 提取表格
    print(f"\n📊 表格提取（共{len(doc.tables)}个表格）:")
    for idx, table in enumerate(doc.tables):
        table_rows = []
        table_info = {
            "index": idx,
            "rows": [],
            "columns": len(table.columns) if table.rows else 0
        }
        
        for row_idx, row in enumerate(table.rows):
            row_cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                row_cells.append(cell_text)
            
            if any(cell for cell in row_cells):
                table_rows.append(" | ".join(row_cells))
                table_info["rows"].append({
                    "row_index": row_idx,
                    "cells": row_cells
                })
        
        if table_rows:
            table_content = "\n".join(table_rows)
            content_parts.append(table_content)
            detailed_info["tables"].append(table_info)
            
            print(f"   表格{idx}:")
            print(f"     行数: {len(table_info['rows'])}")
            print(f"     列数: {table_info['columns']}")
            print(f"     内容预览: {table_content[:200]}{'...' if len(table_content) > 200 else ''}")
    
    # 分析结构
    print(f"\n🔍 结构分析:")
    print(f"   总段落数: {len(detailed_info['paragraphs'])}")
    print(f"   总表格数: {len(detailed_info['tables'])}")
    print(f"   总内容块数: {len(content_parts)}")
    
    # 识别可能的标题和正文
    titles = []
    bodies = []
    for para_info in detailed_info["paragraphs"]:
        if para_info["is_bold"] or para_info["font_size"] and para_info["font_size"] > 12:
            titles.append(para_info)
        else:
            bodies.append(para_info)
    
    print(f"   可能的标题段落: {len(titles)}")
    print(f"   可能的正文段落: {len(bodies)}")
    
    detailed_info["structure"] = {
        "total_paragraphs": len(detailed_info["paragraphs"]),
        "total_tables": len(detailed_info["tables"]),
        "title_paragraphs": len(titles),
        "body_paragraphs": len(bodies)
    }
    
    return {
        "content": "\n\n".join(content_parts),
        "detailed_info": detailed_info
    }

async def test_docx_to_ppt_full_flow():
    """完整流程测试：从docx到PPT（带详细探针）"""
    
    print("\n" + "="*80)
    print("完整流程测试：从Demo文档.docx生成PPT（带详细探针）")
    print("="*80)
    
    # ========== 探针1: 文件检查 ==========
    print("\n" + "="*80)
    print("【探针1】文件检查")
    print("="*80)
    docx_path = Path("Demo文档.docx")
    framework_ppt = Path("demo_filled.pptx")
    
    if not docx_path.exists():
        print(f"❌ 未找到文件: {docx_path}")
        return
    print(f"✅ docx文件存在: {docx_path} ({docx_path.stat().st_size:,} bytes)")
    
    if not framework_ppt.exists():
        print(f"❌ 未找到框架PPT: {framework_ppt}")
        print("   将创建一个新的16:9框架PPT...")
        from create_framework_ppt import create_framework_ppt
        framework_ppt_str = create_framework_ppt()
        framework_ppt = Path(framework_ppt_str)
    print(f"✅ 框架PPT存在: {framework_ppt} ({framework_ppt.stat().st_size:,} bytes)")
    
    # 检查框架PPT的幻灯片数
    from pptx import Presentation
    prs = Presentation(str(framework_ppt))
    print(f"📊 框架PPT信息:")
    print(f"   幻灯片数: {len(prs.slides)}")
    total_placeholders = 0
    for i, slide in enumerate(prs.slides):
        placeholders = [s for s in slide.shapes if s.is_placeholder]
        total_placeholders += len(placeholders)
        print(f"   幻灯片{i}: {len(placeholders)}个占位符")
    print(f"   总占位符数: {total_placeholders}")
    
    # ========== 探针2: 提取docx内容（详细） ==========
    docx_result = extract_docx_content(str(docx_path))
    docx_content = docx_result["content"]
    docx_detailed = docx_result["detailed_info"]
    
    print(f"\n✅ 提取完成")
    print(f"   内容长度: {len(docx_content)} 字符")
    print(f"   段落数: {docx_detailed['structure']['total_paragraphs']}")
    print(f"   表格数: {docx_detailed['structure']['total_tables']}")
    print(f"   内容预览: {docx_content[:200]}...")
    
    # ========== 探针3: 初始化PPT填充器 ==========
    print("\n" + "="*80)
    print("【探针3】初始化PPT填充器（浏览器渲染模式）")
    print("="*80)
    filler = PPTFiller(
        str(framework_ppt),
        use_browser_rendering=True
    )
    print(f"✅ PPT填充器初始化完成")
    print(f"   框架路径: {filler.framework_path}")
    print(f"   浏览器渲染: {filler.use_browser_rendering}")
    print(f"   LLM服务: {'已初始化' if filler.llm_service else '未初始化'}")
    
    # ========== 探针4: 构建用户提示词 ==========
    print("\n" + "="*80)
    print("【探针4】构建生成提示词")
    print("="*80)
    user_prompt = f"""
基于以下文档内容，生成一份完整的PPT演示文稿：

【文档内容】
{docx_content}

【生成要求】
1. 保持文档的核心思想和主要观点
2. 将内容组织成清晰的板块结构
3. 突出关键数据和案例
4. 符合中国商业汇报习惯
5. 使用专业、正式的表达风格
6. 确保内容完整、逻辑清晰
7. 为所有幻灯片生成内容（框架PPT有{len(prs.slides)}张幻灯片）
"""
    print(f"✅ 提示词构建完成")
    print(f"   提示词长度: {len(user_prompt)} 字符")
    print(f"   包含docx内容: {len(docx_content)} 字符")
    
    # ========== 探针5: 执行完整流程 ==========
    print("\n" + "="*80)
    print("【探针5】执行完整流程：生成PPT")
    print("="*80)
    print("   这将执行以下完整流程：")
    print("   1. 提取框架结构（增强解析）")
    print("   2. 人类中心化分析（6层分析）")
    print("   3. 内容生成策略制定")
    print("   4. 智能识别支撑材料（数据点、案例）")
    print("   5. 逐板块内容生成（整合支撑材料）")
    print("     5.1 内容润色（ContentPolisher）")
    print("     5.2 展示策划（PresentationPlanner）")
    print("     5.3 布局规划（LayoutPlanner）【新增】")
    print("   6. HTML生成（基于布局规划，Ant Design规范 + 24栅格系统）【新增】")
    print("   7. 浏览器渲染（Playwright）")
    print("   8. 元素分析和提取（容器、文本）")
    print("   9. 复刻到PPT（坐标映射、24栅格系统）")
    print("   10. 图表生成和整合（如果有数据）")
    print("   11. 最终PPT保存")
    print("")
    
    # 保存日志到文件
    log_file = f"docx_to_ppt_test_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"
    print(f"📝 详细日志将保存到: {log_file}")
    print("")
    
    # 配置日志系统
    setup_logging(log_file)
    logger.info("="*80)
    logger.info("开始完整流程测试")
    logger.info("="*80)
    
    # 【探针】记录关键参数
    logger.info(f"【探针】测试参数:")
    logger.info(f"  - docx文件: {docx_path}")
    logger.info(f"  - 框架PPT: {framework_ppt}")
    logger.info(f"  - 提示词长度: {len(user_prompt)} 字符")
    logger.info(f"  - docx内容长度: {len(docx_content)} 字符")
    logger.info(f"  - 使用增强分析: True")
    logger.info(f"  - 使用浏览器渲染: {filler.use_browser_rendering}")
    logger.info(f"  - LLM服务: {'已初始化' if filler.llm_service else '未初始化'}")
    
    # 检查LLM服务
    if not filler.llm_service:
        logger.error("❌ LLM服务未初始化，无法继续测试")
        print("\n❌ 错误: LLM服务未初始化")
        print("   请确保设置了环境变量: CHAT_MODEL_API_KEY 或 OPENAI_API_KEY")
        return
    
    # 【新增】跳过PPT转换，仅生成HTML
    skip_ppt = True  # 设置为True以跳过HTML到PPT的转换，仅生成HTML文件
    logger.info(f"  - 跳过PPT转换: {skip_ppt} (仅生成HTML)")
    
    output_path = await filler.fill_from_prompt(
        user_prompt,
        output_path="docx_to_ppt_output.pptx",
        use_enhanced_analysis=True,  # 使用增强分析（人类中心化）
        skip_ppt_conversion=skip_ppt  # 跳过HTML到PPT的转换
    )
    
    logger.info("="*80)
    logger.info("完整流程执行完成")
    logger.info("="*80)
    
    # ========== 探针6: 验证输出结果 ==========
    print("\n" + "="*80)
    print("【探针6】验证输出结果")
    print("="*80)
    output_path_obj = Path(output_path)
    
    # 【新增】如果跳过PPT转换，验证HTML文件
    if skip_ppt:
        if output_path_obj.exists() and output_path_obj.is_dir():
            html_files = sorted(output_path_obj.glob("*.html"))
            print(f"✅ HTML输出目录存在: {output_path}")
            print(f"   HTML文件数量: {len(html_files)}")
            for html_file in html_files:
                file_size = html_file.stat().st_size
                print(f"   - {html_file.name}: {file_size:,} bytes ({file_size/1024:.2f} KB)")
            print(f"\n💡 提示: 请打开HTML文件查看效果（在浏览器中打开）")
        else:
            print(f"❌ HTML输出目录不存在: {output_path}")
    else:
        # 验证PPT内容
        if output_path_obj.exists() and output_path_obj.is_file():
            file_size = output_path_obj.stat().st_size
            print(f"✅ 输出文件存在: {output_path}")
            print(f"   文件大小: {file_size:,} bytes ({file_size/1024:.2f} KB)")
            
            # 验证PPT内容
            output_prs = Presentation(str(output_path))
            print(f"📊 输出PPT信息:")
            print(f"   幻灯片数: {len(output_prs.slides)}")
            print(f"   尺寸: {output_prs.slide_width/360000:.2f}cm × {output_prs.slide_height/360000:.2f}cm")
            print(f"   宽高比: {(output_prs.slide_width/output_prs.slide_height):.2f} (16:9 = {16/9:.2f})")
            
            # 检查每张幻灯片的内容
            for i, slide in enumerate(output_prs.slides):
                text_shapes = [s for s in slide.shapes if hasattr(s, 'text') and s.text.strip()]
                image_shapes = [s for s in slide.shapes if hasattr(s, 'image')]
                print(f"   幻灯片{i}:")
                print(f"     文本形状: {len(text_shapes)}")
                print(f"     图片形状: {len(image_shapes)}")
                if text_shapes:
                    print(f"     文本预览: {text_shapes[0].text[:50]}...")
        else:
            print(f"❌ 输出文件不存在: {output_path}")
    
    # ========== 探针7: 分析日志 ==========
    print("\n" + "="*80)
    print("【探针7】分析详细日志")
    print("="*80)
    analyze_log_file(log_file)
    
    # ========== 最终总结 ==========
    print("\n" + "="*80)
    print("✅ 完整流程测试完成！")
    print("="*80)
    print(f"\n📁 输出文件: {output_path}")
    print(f"📝 详细日志: {log_file}")
    
    print("\n" + "="*80)
    print("💡 提示: 请打开生成的PPT文件查看效果")
    print("="*80)

def analyze_log_file(log_file: str):
    """分析日志文件，提取关键信息"""
    log_path = Path(log_file)
    if not log_path.exists():
        print(f"❌ 日志文件不存在: {log_file}")
        return
    
    print(f"📊 分析日志文件: {log_file}")
    print(f"   文件大小: {log_path.stat().st_size:,} bytes")
    
    # 读取日志内容
    with open(log_file, 'r', encoding='utf-8') as f:
        log_lines = f.readlines()
    
    print(f"   总行数: {len(log_lines)}")
    
    # 分析关键阶段
    stages = {
        "人类中心化分析": 0,
        "内容生成策略": 0,
        "支撑材料识别": 0,
        "内容润色": 0,
        "展示策划": 0,
        "布局规划": 0,
        "HTML生成": 0,
        "浏览器渲染": 0,
        "PPT复刻": 0,
        "错误": 0,
        "警告": 0
    }
    
    # 统计各阶段出现次数
    for line in log_lines:
        line_lower = line.lower()
        if "人类中心化分析" in line or "human-centered" in line_lower:
            stages["人类中心化分析"] += 1
        if "内容生成策略" in line or "content strategy" in line_lower:
            stages["内容生成策略"] += 1
        if "支撑材料" in line or "supporting materials" in line_lower:
            stages["支撑材料识别"] += 1
        if "内容润色" in line or "polish" in line_lower or "润色" in line:
            stages["内容润色"] += 1
        if "展示策划" in line or "presentation plan" in line_lower or "展示策划" in line:
            stages["展示策划"] += 1
        if "布局规划" in line or "layout plan" in line_lower or "布局规划" in line:
            stages["布局规划"] += 1
        if "html生成" in line or "generate.*html" in line_lower or "generate_from_layout_plan" in line:
            stages["HTML生成"] += 1
        if "浏览器渲染" in line or "browser render" in line_lower:
            stages["浏览器渲染"] += 1
        if "复刻" in line or "replicate" in line_lower:
            stages["PPT复刻"] += 1
        if "error" in line_lower or "❌" in line or "失败" in line:
            stages["错误"] += 1
        if "warning" in line_lower or "⚠️" in line or "警告" in line:
            stages["警告"] += 1
    
    print(f"\n📈 各阶段统计:")
    for stage, count in stages.items():
        if count > 0:
            print(f"   {stage}: {count} 次")
    
    # 查找关键信息
    print(f"\n🔍 关键信息提取:")
    
    # 查找润色结果
    polished_count = 0
    for i, line in enumerate(log_lines):
        if "润色完成" in line or "polish.*完成" in line.lower():
            polished_count += 1
            if polished_count <= 3:  # 只显示前3个
                print(f"   ✅ 润色完成: {line.strip()}")
    
    # 查找布局规划结果
    layout_count = 0
    for i, line in enumerate(log_lines):
        if "布局规划完成" in line or "layout.*plan.*完成" in line.lower():
            layout_count += 1
            if layout_count <= 3:  # 只显示前3个
                print(f"   ✅ 布局规划完成: {line.strip()}")
    
    # 查找HTML生成方式
    html_method = None
    for line in log_lines:
        if "使用布局规划生成HTML" in line:
            html_method = "布局规划方式"
            break
        elif "使用内容映射生成HTML" in line:
            html_method = "内容映射方式"
            break
    
    if html_method:
        print(f"   📄 HTML生成方式: {html_method}")
    
    # 查找错误和警告
    errors = []
    warnings = []
    for i, line in enumerate(log_lines):
        if "error" in line.lower() or "❌" in line or "失败" in line:
            errors.append((i+1, line.strip()[:100]))
        if "warning" in line.lower() or "⚠️" in line or "警告" in line:
            warnings.append((i+1, line.strip()[:100]))
    
    if errors:
        print(f"\n❌ 发现 {len(errors)} 个错误:")
        for line_num, error_msg in errors[:5]:  # 只显示前5个
            print(f"   行{line_num}: {error_msg}")
    
    if warnings:
        print(f"\n⚠️ 发现 {len(warnings)} 个警告:")
        for line_num, warn_msg in warnings[:5]:  # 只显示前5个
            print(f"   行{line_num}: {warn_msg}")
    
    if not errors and not warnings:
        print(f"\n✅ 未发现错误或警告")

if __name__ == "__main__":
    asyncio.run(test_docx_to_ppt_full_flow())

