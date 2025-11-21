"""
PPT生成器 - 核心组装功能
从FabricatorAgent中抽离的PPT组装核心功能，不依赖数据库和知识库
支持可选的Vinci图表生成集成
"""

import json
import os
import re
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path
from typing import Any, Optional, Union

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt
from webcolors import hex_to_rgb

# 导入Ant Design主题
from ant_design_theme import ant_design_theme


# 辅助类，用于管理布局计算（使用统一单位）
class Box:
    def __init__(self, left, top, width, height):
        # 统一转换为 Cm 类型，确保运算兼容性
        def _to_cm(value):
            if isinstance(value, Cm):
                return value
            elif isinstance(value, (Pt, Emu)):
                return Cm(float(value) / 360000)  # Emu to cm
            elif isinstance(value, (int, float)):
                return Cm(value)
            else:
                return Cm(0)

        self.left = _to_cm(left)
        self.top = _to_cm(top)
        self.width = _to_cm(width)
        self.height = _to_cm(height)

    def _cm_to_float(self, cm_value):
        """将 Cm 对象转换为厘米数值"""
        if isinstance(cm_value, Cm):
            return float(cm_value) / 360000
        elif isinstance(cm_value, (int, float)):
            return float(cm_value)
        else:
            return 0.0


class PPTGenerator:
    """
    独立的PPT生成器
    根据VML计划和内容映射生成PPT文件
    支持可选的Vinci图表生成集成
    """

    def __init__(
        self,
        output_dir: Optional[Union[str, Path]] = None,
        vinci_integration: Optional[Any] = None
    ):
        """
        初始化PPT生成器
        
        Args:
            output_dir: PPT输出目录，默认为当前目录下的 ppt_outputs 文件夹
            vinci_integration: 可选的Vinci集成实例，用于生成图表
        """
        self.codename = "PPTGenerator"
        
        if output_dir is None:
            output_dir = Path.cwd() / "ppt_outputs"
        elif isinstance(output_dir, str):
            output_dir = Path(output_dir)
        
        self.PPT_OUTPUT_DIR = output_dir
        self.PPT_OUTPUT_DIR.mkdir(exist_ok=True, parents=True)
        
        self._vinci_integration = vinci_integration
        if vinci_integration:
            logger.info(f"--- [{self.codename}]: Initialized with Vinci integration")
        else:
            logger.info(f"--- [{self.codename}]: Initialized without Vinci integration (chart generation disabled)")
        
        logger.info(f"--- [{self.codename}]: Output directory: {self.PPT_OUTPUT_DIR}")

    def _parse_unit(self, value_str: str, base_cm: float = 0, base_pt: float = 0) -> Union[Cm, Pt, Emu]:
        """解析并转换各种单位 (cm, px, pt, %) 到 python-pptx 的单位类型"""
        if isinstance(value_str, (int, float)):
            return Cm(value_str)

        value_str = str(value_str).lower().strip()

        try:
            if value_str.endswith('cm'):
                return Cm(float(value_str[:-2]))
            elif value_str.endswith('px'):
                px_value = float(value_str[:-2])
                return Cm(px_value * 2.54 / 96)
            elif value_str.endswith('pt'):
                return Pt(float(value_str[:-2]))
            elif value_str.endswith('%'):
                if base_cm > 0:
                    return Cm(base_cm * (float(value_str[:-1]) / 100.0))
                elif base_pt > 0:
                    return Pt(base_pt * (float(value_str[:-1]) / 100.0))
                return Cm(0)
            else:
                return Cm(float(value_str))
        except (ValueError, TypeError):
            logger.warning(f"--- [{self.codename}] Could not parse unit from '{value_str}'. Defaulting to 0.")
            return Cm(0)

    def _parse_color(self, color_str: str) -> Optional[dict[str, Any]]:
        """解析颜色字符串 (hex 或 rgba) 并返回包含 RGB 和 alpha 的字典"""
        if not color_str:
            return None

        try:
            # 处理 rgba 格式
            if color_str.startswith('rgba'):
                parts = [p.strip() for p in color_str[5:-1].split(',')]
                if len(parts) >= 3:
                    r, g, b = int(parts[0]), int(parts[1]), int(parts[2])
                    a = float(parts[3]) if len(parts) > 3 else 1.0
                    return {"rgb": RGBColor(r, g, b), "alpha": a}

            # 处理 rgb 格式（无 alpha）
            if color_str.startswith('rgb'):
                parts = [p.strip() for p in color_str[4:-1].split(',')]
                if len(parts) >= 3:
                    r, g, b = int(parts[0]), int(parts[1]), int(parts[2])
                    return {"rgb": RGBColor(r, g, b), "alpha": 1.0}

            # 处理 hex 格式
            if color_str.startswith('#'):
                rgb = hex_to_rgb(color_str)
                return {"rgb": RGBColor(rgb.red, rgb.green, rgb.blue), "alpha": 1.0}

            # 处理不带 # 的 hex
            if len(color_str) == 6 and all(c in '0123456789abcdef' for c in color_str.lower()):
                rgb = hex_to_rgb('#' + color_str)
                return {"rgb": RGBColor(rgb.red, rgb.green, rgb.blue), "alpha": 1.0}

        except Exception as e:
            logger.warning(f"--- [{self.codename}] Could not parse color '{color_str}'. Error: {e}")

        return None

    def _parse_shadow(self, shadow_str: str) -> Optional[dict[str, Any]]:
        """解析 CSS box-shadow 字符串，转换为 python-pptx 阴影参数"""
        if not shadow_str or shadow_str.lower() == 'none':
            return None

        try:
            first_shadow = shadow_str.split(',')[0].strip()
            color_match = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', first_shadow)
            if not color_match:
                return None

            r, g, b, a = color_match.groups()
            a = float(a) if a is not None else 1.0

            offset_matches = re.findall(r'(-?\d+(?:\.\d+)?)px', first_shadow)
            if len(offset_matches) < 2:
                return None

            offset_x = float(offset_matches[0])
            offset_y = float(offset_matches[1])
            blur_radius = float(offset_matches[2]) if len(offset_matches) > 2 else 0

            distance = Pt(abs(offset_y))
            direction = 270 if offset_y >= 0 else 90

            return {
                "color": RGBColor(int(r), int(g), int(b)),
                "transparency": 1.0 - a,
                "blur_radius": Pt(blur_radius),
                "distance": distance,
                "direction": direction
            }
        except Exception as e:
            logger.warning(f"--- [{self.codename}] Could not parse shadow '{shadow_str}'. Error: {e}")
        return None

    def _parse_vml(self, vml_code: str):
        """解析 VML 代码，如果失败则尝试修复或返回安全的 fallback"""
        try:
            return ET.fromstring(vml_code)
        except ET.ParseError as e:
            try:
                from xml.sax.saxutils import escape

                def escape_text_content(match):
                    attr_name = match.group(1)
                    content = match.group(2)
                    escaped_content = escape(content)
                    return f'{attr_name}="{escaped_content}"'

                fixed_vml = re.sub(r'(text)="([^"]*)"', escape_text_content, vml_code)

                if fixed_vml != vml_code:
                    logger.debug("--- [PPTGenerator] VML parse failed, attempting auto-fix...")
                    try:
                        result = ET.fromstring(fixed_vml)
                        logger.debug("--- [PPTGenerator] ✅ VML auto-fixed and parsed successfully!")
                        return result
                    except ET.ParseError:
                        logger.warning(f"--- [PPTGenerator] Auto-fix failed. Using fallback.")
            except Exception as fix_error:
                logger.warning(f"--- [PPTGenerator] Failed to fix VML: {fix_error}")

            logger.error(f"VML Parse Error: {e}. VML Code:\n{vml_code}")

            # 生成安全的 fallback VML
            from xml.sax.saxutils import escape
            error_msg = escape(str(e))
            try:
                fallback_vml = f'<Slide padding="1.5cm"><VStack><TextBox style="title" text="VML解析失败"/><TextBox style="body" text="错误: {error_msg}"/></VStack></Slide>'
                return ET.fromstring(fallback_vml)
            except ET.ParseError:
                logger.critical("Fallback VML also failed. Using minimal safe VML.")
                return ET.fromstring('<Slide padding="1.5cm"><VStack><TextBox style="title" text="VML解析失败"/></VStack></Slide>')

    def _render_element(self, slide, element, box: Box, content_map: dict):
        """渲染VML元素到PPT幻灯片"""
        tag = element.tag.lower()
        attrs = {k: v for k, v in element.attrib.items()}

        # 容器元素递归渲染
        if tag in ["vstack", "hstack", "stack"]:
            children = list(element)
            if not children:
                return

            direction = attrs.get("direction", "vertical").lower()
            if tag == "vstack":
                direction = "vertical"
            elif tag == "hstack":
                direction = "horizontal"

            gap_str = attrs.get("gap", "0.5cm")
            base_dim = box.width if direction == 'horizontal' else box.height
            base_dim_cm = box._cm_to_float(base_dim) if isinstance(base_dim, Cm) else float(base_dim)
            gap = self._parse_unit(gap_str, base_cm=base_dim_cm)

            if len(children) > 1:
                if isinstance(gap, Cm):
                    gap_cm = box._cm_to_float(gap)
                    total_gap_cm = gap_cm * (len(children) - 1)
                    total_gap = Cm(total_gap_cm)
                else:
                    total_gap = Cm(float(gap) * (len(children) - 1))
            else:
                total_gap = Cm(0)

            if direction == "vertical":
                total_gap_cm = box._cm_to_float(total_gap)
                box_height_cm = box._cm_to_float(box.height)
                available_height_cm = box_height_cm - total_gap_cm if len(children) > 1 else box_height_cm
                child_height_cm = available_height_cm / len(children) if children else 0
                child_height = Cm(child_height_cm)

                current_top = box.top
                gap_cm = box._cm_to_float(gap) if isinstance(gap, Cm) else float(gap)
                for child in children:
                    child_box = Box(box.left, current_top, box.width, child_height)
                    self._render_element(slide, child, child_box, content_map)
                    current_top_cm = box._cm_to_float(current_top) + child_height_cm + gap_cm
                    current_top = Cm(current_top_cm)
            else:  # horizontal
                total_gap_cm = box._cm_to_float(total_gap)
                box_width_cm = box._cm_to_float(box.width)
                available_width_cm = box_width_cm - total_gap_cm if len(children) > 1 else box_width_cm
                child_width_cm = available_width_cm / len(children) if children else 0
                child_width = Cm(child_width_cm)

                current_left = box.left
                gap_cm = box._cm_to_float(gap) if isinstance(gap, Cm) else float(gap)
                for child in children:
                    child_box = Box(current_left, box.top, child_width, box.height)
                    self._render_element(slide, child, child_box, content_map)
                    current_left_cm = box._cm_to_float(current_left) + child_width_cm + gap_cm
                    current_left = Cm(current_left_cm)
            return

        # 判断是否需要创建带样式的容器
        is_styled_container = (
            'background' in attrs or
            'border' in attrs or
            'shadows' in attrs or
            'shadow' in attrs or
            'borderRadius' in attrs or
            'borderradius' in attrs
        )

        # 1. 处理纯文本内容（无样式容器）
        if tag == "textbox" and not is_styled_container:
            text_ref = attrs.get("ref", "")
            if not text_ref:
                text = attrs.get("text", "内容缺失")
                logger.debug(f"--- [PPTGenerator] Rendering pure TextBox (fallback): text='{text[:50]}...'")
            else:
                text = content_map.get(text_ref, f"!!REF_NOT_FOUND: {text_ref}!!")
                logger.debug(f"--- [PPTGenerator] Rendering pure TextBox: ref='{text_ref}'")

            txBox = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
            tf = txBox.text_frame
            tf.clear()
            tf.word_wrap = True

            justify = attrs.get("justify", "top").lower()
            if justify == "center":
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            elif justify == "bottom":
                tf.vertical_anchor = MSO_ANCHOR.BOTTOM
            else:
                tf.vertical_anchor = MSO_ANCHOR.TOP

            p = tf.paragraphs[0]
            align = attrs.get("align", "left").lower()
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            p.text = str(text)

            try:
                if hasattr(txBox, 'line'):
                    txBox.line.fill.background()
            except Exception:
                pass

            try:
                if hasattr(txBox, 'fill'):
                    txBox.fill.background()
            except Exception:
                pass

            if not p.runs:
                run = p.add_run()
            else:
                run = p.runs[0]

            font = run.font
            # 使用Ant Design字体系统
            try:
                # Ant Design字体栈（优先系统字体）
                font.name = "Segoe UI"  # Windows系统字体
            except Exception:
                try:
                    font.name = "Helvetica Neue"  # macOS系统字体
                except Exception:
                    try:
                        font.name = "微软雅黑"  # 中文字体fallback
                    except Exception:
                        font.name = "Arial"  # 最终fallback

            if 'color' in attrs:
                font_color_info = self._parse_color(attrs['color'])
                if font_color_info:
                    font.color.rgb = font_color_info['rgb']

            if 'fontSize' in attrs:
                font_size = self._parse_unit(attrs['fontSize'], base_pt=100)
                if isinstance(font_size, Pt):
                    font.size = font_size
                else:
                    font.size = Pt(float(font_size) * 28.35)

            if 'fontWeight' in attrs:
                try:
                    weight = int(attrs['fontWeight'])
                    font.bold = weight >= 600
                except (ValueError, TypeError):
                    font.bold = attrs.get('fontWeight', '').lower() in ['bold', '700', '600']

            if 'fontSize' not in attrs and 'fontWeight' not in attrs:
                style_key = attrs.get("style", "body")
                if style_key == "title":
                    font.size = Pt(ant_design_theme.get_font_size_pt('h1'))
                    font.bold = True
                elif style_key == "subtitle":
                    font.size = Pt(ant_design_theme.get_font_size_pt('h3'))
                    font.bold = False
                else:
                    font.size = Pt(ant_design_theme.get_font_size_pt('base'))
                    font.bold = False

            return

        # 2. 处理纯图片内容（无样式容器）
        elif tag in ["imagebox", "image"] and not is_styled_container:
            img_ref = attrs.get("ref", "")
            img_path = content_map.get(img_ref)

            logger.info(f"--- [PPTGenerator] Rendering pure ImageBox: ref='{img_ref}', found_path='{img_path}'")

            if not img_path or not os.path.exists(img_path):
                logger.error(f"--- [PPTGenerator] CRITICAL: Image path not found: {img_path}")
                txBox = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
                txBox.text_frame.text = f"Image Load Error!\nRef: {img_ref}"
            else:
                width_str = attrs.get("width", "100%")
                height_str = attrs.get("height", "100%")

                box_width_cm = box._cm_to_float(box.width)
                box_height_cm = box._cm_to_float(box.height)
                img_width = self._parse_unit(width_str, base_cm=box_width_cm)
                img_height = self._parse_unit(height_str, base_cm=box_height_cm)

                try:
                    slide.shapes.add_picture(str(img_path), box.left, box.top, img_width, img_height)
                except Exception as e:
                    logger.error(f"--- [PPTGenerator] Failed to add picture: {e}", exc_info=True)
                    txBox = slide.shapes.add_textbox(box.left, box.top, box.width, box.height)
                    txBox.text_frame.text = f"Image Error: {img_ref}\n{str(e)}"

            return

        # 3. 处理带样式的容器
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if ('borderRadius' in attrs or 'borderradius' in attrs) else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, box.left, box.top, box.width, box.height)

        if 'borderRadius' in attrs or 'borderradius' in attrs:
            try:
                radius_str = attrs.get('borderRadius') or attrs.get('borderradius', '')
                radius_px = float(re.sub(r'[a-zA-Z%]', '', str(radius_str)))
                box_height_cm = box._cm_to_float(box.height)
                adjustment = min(0.5, radius_px / 32.0)
                if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE and len(shape.adjustments) > 0:
                    shape.adjustments[0] = adjustment
            except Exception as e:
                logger.warning(f"--- [PPTGenerator] Failed to apply custom borderRadius: {e}")

        fill = shape.fill
        line = shape.line
        fill.background()
        line.fill.background()

        if 'background' in attrs:
            color_info = self._parse_color(attrs['background'])
            if color_info:
                fill.solid()
                fill.fore_color.rgb = color_info['rgb']
                if color_info['alpha'] < 1.0:
                    fill.transparency = 1.0 - color_info['alpha']

        if 'shadows' in attrs or 'shadow' in attrs:
            shadow_params = self._parse_shadow(attrs.get('shadows') or attrs.get('shadow'))
            if shadow_params:
                try:
                    shadow = shape.shadow
                    shadow.inherit = False
                    shadow.style = 'outer'
                    shadow.blur_radius = shadow_params['blur_radius']
                    shadow.distance = shadow_params['distance']
                    shadow.direction = shadow_params['direction']
                    shadow.color.rgb = shadow_params['color']
                    shadow.transparency = shadow_params['transparency']
                except Exception as e:
                    logger.warning(f"--- [PPTGenerator] Could not apply shadow. Error: {e}")

        # 4. 处理文本内容（在带样式的容器中）
        if tag == "textbox":
            text_ref = attrs.get("ref", "")
            if not text_ref:
                text = attrs.get("text", "内容缺失")
                logger.debug(f"--- [PPTGenerator] Rendering TextBox (fallback): text='{text[:50]}...'")
            else:
                text = content_map.get(text_ref, f"!!REF_NOT_FOUND: {text_ref}!!")
                logger.debug(f"--- [PPTGenerator] Rendering TextBox: ref='{text_ref}'")

            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            justify = attrs.get("justify", "top").lower()
            if justify == "center":
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            elif justify == "bottom":
                tf.vertical_anchor = MSO_ANCHOR.BOTTOM
            else:
                tf.vertical_anchor = MSO_ANCHOR.TOP

            p = tf.paragraphs[0]
            align = attrs.get("align", "left").lower()
            if align == "center":
                p.alignment = PP_ALIGN.CENTER
            elif align == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            p.text = str(text)

            font = p.font
            # 使用Ant Design字体系统
            try:
                font.name = "Segoe UI"
            except Exception:
                try:
                    font.name = "Helvetica Neue"
                except Exception:
                    try:
                        font.name = "微软雅黑"
                    except Exception:
                        font.name = "Arial"

            if 'color' in attrs:
                font_color_info = self._parse_color(attrs['color'])
                if font_color_info:
                    font.color.rgb = font_color_info['rgb']

            if 'fontSize' in attrs:
                font_size = self._parse_unit(attrs['fontSize'], base_pt=100)
                if isinstance(font_size, Pt):
                    font.size = font_size
                else:
                    font.size = Pt(float(font_size) * 28.35)

            if 'fontWeight' in attrs:
                try:
                    weight = int(attrs['fontWeight'])
                    font.bold = weight >= 600
                except (ValueError, TypeError):
                    font.bold = attrs.get('fontWeight', '').lower() in ['bold', '700', '600']

            if 'fontSize' not in attrs and 'fontWeight' not in attrs:
                style_key = attrs.get("style", "body")
                if style_key == "title":
                    font.size = Pt(ant_design_theme.get_font_size_pt('h1'))
                    font.bold = True
                elif style_key == "subtitle":
                    font.size = Pt(ant_design_theme.get_font_size_pt('h3'))
                    font.bold = False
                else:
                    font.size = Pt(ant_design_theme.get_font_size_pt('base'))
                    font.bold = False

        # 5. 处理图片内容（在带样式的容器中）
        elif tag in ["imagebox", "image"]:
            img_ref = attrs.get("ref", "")
            img_path = content_map.get(img_ref)

            logger.info(f"--- [PPTGenerator] Rendering ImageBox/Image in styled container: ref='{img_ref}'")

            if not img_path or not os.path.exists(img_path):
                logger.error(f"--- [PPTGenerator] CRITICAL: Image path not found: {img_path}")
                shape.text_frame.text = f"Image Load Error!\nRef: {img_ref}"
            else:
                width_str = attrs.get("width", "100%")
                height_str = attrs.get("height", "100%")

                box_width_cm = box._cm_to_float(box.width)
                box_height_cm = box._cm_to_float(box.height)
                img_width = self._parse_unit(width_str, base_cm=box_width_cm)
                img_height = self._parse_unit(height_str, base_cm=box_height_cm)

                try:
                    fill.picture(str(img_path))
                    fill.tile = False
                except (AttributeError, Exception) as e:
                    logger.warning(f"--- [PPTGenerator] fill.picture() not available, using add_picture: {e}")
                    try:
                        pic_shape = slide.shapes.add_picture(str(img_path), box.left, box.top, img_width, img_height)
                        try:
                            shape.fill.background()
                            shape.line.fill.background()
                        except Exception:
                            pass
                    except Exception as e2:
                        logger.error(f"--- [PPTGenerator] Failed to add picture: {e2}", exc_info=True)
                        try:
                            shape.text_frame.text = f"Image Error: {img_ref}\n{str(e2)}"
                        except:
                            pass

    def assemble_slide_from_vml(self, slide, vml_code: str, content_map: dict):
        """从VML代码组装单张幻灯片"""
        root_element = self._parse_vml(vml_code)
        attrs = root_element.attrib

        if root_element.tag.lower() == "slide":
            try:
                slide_width = slide.part.package.presentation_part.presentation.slide_width
                slide_height = slide.part.package.presentation_part.presentation.slide_height
                if isinstance(slide_width, Cm):
                    slide_width_cm = float(slide_width)
                else:
                    slide_width_cm = float(slide_width) / 360000

                if isinstance(slide_height, Cm):
                    slide_height_cm = float(slide_height)
                else:
                    slide_height_cm = float(slide_height) / 360000
            except Exception as e:
                logger.warning(f"--- [PPTGenerator] Could not get slide dimensions: {e}. Using defaults.")
                slide_width_cm = 33.867
                slide_height_cm = 19.05

            # 使用Ant Design间距系统（默认padding）
            default_padding_cm = ant_design_theme.get_spacing_cm('lg')
            padding_str = attrs.get("padding", f"{default_padding_cm:.2f}cm")
            padding = self._parse_unit(padding_str, base_cm=slide_width_cm)
            if isinstance(padding, Cm):
                padding_cm = float(padding) / 360000
            else:
                padding_cm = float(padding)

            inner_box = Box(
                left=Cm(padding_cm),
                top=Cm(padding_cm),
                width=Cm(slide_width_cm - (padding_cm * 2)),
                height=Cm(slide_height_cm - (padding_cm * 2))
            )

            if 'background' in attrs:
                try:
                    color = self._parse_color(attrs['background'])
                    if color:
                        slide.background.fill.solid()
                        slide.background.fill.fore_color.rgb = color['rgb']
                except Exception as e:
                    logger.warning(f"--- [PPTGenerator] Could not apply slide background color. Error: {e}")

            for child in list(root_element):
                self._render_element(slide, child, inner_box, content_map)

    async def generate_ppt(
        self,
        project_name: str,
        vml_plan: list[dict],
        content_map: dict,
        template_path: Optional[Union[str, Path]] = None,
        chart_insights: Optional[list[dict]] = None
    ) -> dict:
        """
        根据VML计划生成PPT文件
        
        Args:
            project_name: 项目名称（用于文件名）
            vml_plan: VML计划列表，每个元素包含 'vml_code' 字段
            content_map: 内容映射字典，键是ref名称，值是实际内容（文本或图片路径）
            template_path: 可选的PPT模板路径
            chart_insights: 可选的图表洞察列表，如果提供且Vinci集成可用，会自动生成图表
            
        Returns:
            包含 'file_path' 的字典，如果失败则包含 'error'
        """
        logger.info(f"--- [{self.codename}]: Generating PPT from VML plan...")
        logger.info(f"VML Plan ({len(vml_plan)} slides)")
        logger.info(f"Content Map ({len(content_map)} items)")
        
        # 如果提供了图表洞察且Vinci集成可用，先生成图表
        if chart_insights and self._vinci_integration:
            logger.info(f"--- [{self.codename}]: Generating {len(chart_insights)} charts via Vinci...")
            try:
                chart_paths = await self._vinci_integration.generate_charts_from_insights(
                    chart_insights,
                    project_id=project_name
                )
                # 将生成的图表路径添加到content_map
                for insight_id, chart_path in chart_paths.items():
                    # 如果content_map中已经有这个ref，更新它；否则添加
                    if insight_id in content_map:
                        logger.info(f"--- [{self.codename}]: Updated chart path for '{insight_id}': {chart_path}")
                    content_map[insight_id] = chart_path
                logger.success(f"--- [{self.codename}]: Generated {len(chart_paths)} charts")
            except Exception as e:
                logger.error(f"--- [{self.codename}]: Failed to generate charts: {e}", exc_info=True)
                # 继续生成PPT，即使图表生成失败
        elif chart_insights and not self._vinci_integration:
            logger.warning(
                f"--- [{self.codename}]: Chart insights provided but Vinci integration not available. "
                "Charts will not be generated."
            )

        try:
            # 加载模板或创建空白演示文稿
            if template_path and Path(template_path).exists():
                prs = Presentation(str(template_path))
                logger.info(f"--- [{self.codename}]: Loaded template from: {template_path}")
            else:
                prs = Presentation()
                prs.slide_width = Cm(33.867)  # 16:9
                prs.slide_height = Cm(19.05)
                logger.info(f"--- [{self.codename}]: Using blank presentation")

            # 选择空白布局
            num_layouts = len(prs.slide_layouts)
            blank_layout = None
            min_placeholders = float('inf')

            for layout_idx in range(num_layouts):
                try:
                    test_layout = prs.slide_layouts[layout_idx]
                    placeholder_count = len(test_layout.placeholders)
                    if placeholder_count < min_placeholders:
                        min_placeholders = placeholder_count
                        blank_layout = test_layout
                    if placeholder_count == 0:
                        break
                except Exception:
                    continue

            if blank_layout is None:
                blank_layout = prs.slide_layouts[0]
                logger.warning(f"--- [{self.codename}] Could not find blank layout, using first layout")
            else:
                logger.info(f"--- [{self.codename}] Selected blank layout: {blank_layout.name} ({min_placeholders} placeholders)")

            # 处理每张幻灯片
            for i, slide_vml_data in enumerate(vml_plan):
                vml_code = slide_vml_data.get("vml_code", "")

                try:
                    slide = prs.slides.add_slide(blank_layout)
                except Exception as layout_error:
                    logger.warning(f"--- [{self.codename}] Failed to add slide: {layout_error}")
                    slide = prs.slides.add_slide(prs.slide_layouts[0])

                # 清除占位符
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        try:
                            if hasattr(shape, 'text_frame'):
                                shape.text_frame.clear()
                            if hasattr(shape, 'fill'):
                                try:
                                    shape.fill.background()
                                except:
                                    pass
                            if hasattr(shape, 'line'):
                                try:
                                    shape.line.fill.background()
                                except:
                                    pass
                        except Exception:
                            pass

                try:
                    self.assemble_slide_from_vml(slide, vml_code, content_map)
                except Exception as slide_error:
                    logger.error(f"--- [{self.codename}] Failed to render slide {i+1}: {slide_error}", exc_info=True)
                    try:
                        error_textbox = slide.shapes.add_textbox(Cm(2), Cm(2), Cm(10), Cm(2))
                        error_textbox.text_frame.text = f"幻灯片 {i+1} 渲染失败: {str(slide_error)[:100]}"
                    except:
                        pass

            # 保存文件
            safe_project_name = "".join(c for c in project_name if c.isalnum() or c in " -_").rstrip()
            timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
            file_name = f"{safe_project_name}-{timestamp}.pptx"
            file_path = self.PPT_OUTPUT_DIR / file_name

            try:
                prs.save(str(file_path))

                if not file_path.exists():
                    raise FileNotFoundError(f"PPT file was not created: {file_path}")

                file_size = file_path.stat().st_size
                if file_size < 1000:
                    logger.warning(f"--- [{self.codename}] Generated PPT file is suspiciously small: {file_size} bytes")

                logger.success(f"--- [{self.codename}]: PPT file generated and saved to: {file_path} ({file_size} bytes)")
                return {"file_path": str(file_path), "file_size": file_size}
            except Exception as save_error:
                logger.error(f"--- [{self.codename}] Failed to save PPT file: {save_error}", exc_info=True)
                raise

        except Exception as e:
            error_msg = f"Failed to generate PPT: {e}"
            logger.error(f"--- [{self.codename}]: {error_msg}", exc_info=True)
            return {"error": error_msg}

