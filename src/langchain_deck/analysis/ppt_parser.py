"""
PPT框架解析器
读取现有PPT文件，提取结构和布局信息
"""

from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from loguru import logger


class PPTParser:
    """
    PPT框架解析器
    从现有PPT文件中提取结构和内容信息
    """
    
    def __init__(self, ppt_path: str):
        """
        初始化PPT解析器
        
        Args:
            ppt_path: PPT文件路径
        """
        self.ppt_path = Path(ppt_path)
        if not self.ppt_path.exists():
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")
        
        self.prs = Presentation(str(self.ppt_path))
        logger.info(f"--- [PPTParser]: Loaded PPT: {self.ppt_path}")
    
    def extract_structure(self) -> Dict[str, Any]:
        """
        提取PPT的结构信息
        
        Returns:
            包含幻灯片结构信息的字典
        """
        structure = {
            "slides": [],
            "slide_count": len(self.prs.slides),
            "slide_width": float(self.prs.slide_width) / 360000,  # 转换为cm
            "slide_height": float(self.prs.slide_height) / 360000
        }
        
        for idx, slide in enumerate(self.prs.slides):
            slide_info = {
                "slide_index": idx,
                "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
                "shapes": [],
                "placeholders": [],
                "text_content": []
            }
            
            # 提取所有形状信息
            for shape in slide.shapes:
                shape_info = self._extract_shape_info(shape, idx)
                if shape_info:
                    slide_info["shapes"].append(shape_info)
                    
                    # 如果是占位符，单独记录
                    if shape.is_placeholder:
                        slide_info["placeholders"].append(shape_info)
                    
                    # 如果有文本内容，记录
                    if shape_info.get("text"):
                        slide_info["text_content"].append({
                            "type": shape_info["type"],
                            "text": shape_info["text"],
                            "placeholder_id": shape_info.get("placeholder_id")
                        })
            
            structure["slides"].append(slide_info)
        
        logger.info(f"--- [PPTParser]: Extracted structure from {len(structure['slides'])} slides")
        return structure
    
    def _extract_shape_info(self, shape, slide_index: int) -> Optional[Dict[str, Any]]:
        """
        提取单个形状的信息
        
        Args:
            shape: PPT形状对象
            slide_index: 幻灯片索引
            
        Returns:
            形状信息字典
        """
        try:
            shape_info = {
                "type": self._get_shape_type(shape),
                "shape_id": shape.shape_id,
                "left": float(shape.left) / 360000,  # 转换为cm
                "top": float(shape.top) / 360000,
                "width": float(shape.width) / 360000,
                "height": float(shape.height) / 360000,
                "is_placeholder": shape.is_placeholder
            }
            
            # 如果是占位符，记录占位符信息
            if shape.is_placeholder:
                try:
                    shape_info["placeholder_id"] = shape.placeholder_format.idx
                    shape_info["placeholder_type"] = str(shape.placeholder_format.type)
                except:
                    pass
            
            # 提取文本内容
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    shape_info["text"] = text
                    shape_info["has_text"] = True
                else:
                    shape_info["has_text"] = False
                    shape_info["text"] = ""
            elif hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    shape_info["text"] = text
                    shape_info["has_text"] = True
                else:
                    shape_info["has_text"] = False
                    shape_info["text"] = ""
            else:
                shape_info["has_text"] = False
                shape_info["text"] = ""
            
            # 提取图片信息
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    shape_info["image_path"] = shape.image.filename if hasattr(shape.image, 'filename') else None
                except:
                    pass
            
            return shape_info
            
        except Exception as e:
            logger.warning(f"--- [PPTParser]: Failed to extract shape info: {e}")
            return None
    
    def _get_shape_type(self, shape) -> str:
        """获取形状类型名称"""
        try:
            shape_type = shape.shape_type
            type_names = {
                MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
                MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
                MSO_SHAPE_TYPE.PICTURE: "picture",
                MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
                MSO_SHAPE_TYPE.GROUP: "group",
                MSO_SHAPE_TYPE.TABLE: "table",
                MSO_SHAPE_TYPE.MEDIA: "media"
            }
            return type_names.get(shape_type, "unknown")
        except:
            return "unknown"
    
    def extract_text_summary(self) -> str:
        """
        提取PPT的文本摘要，用于LLM理解框架内容
        
        Returns:
            文本摘要字符串
        """
        summary_parts = []
        summary_parts.append(f"PPT框架文档包含 {len(self.prs.slides)} 张幻灯片。\n")
        
        for idx, slide in enumerate(self.prs.slides):
            summary_parts.append(f"\n幻灯片 {idx + 1}:")
            
            # 提取所有文本内容
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
                elif hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            
            if texts:
                summary_parts.append(f"  内容: {' | '.join(texts)}")
            else:
                summary_parts.append(f"  内容: (空白占位符)")
            
            # 记录占位符信息
            placeholders = [s for s in slide.shapes if s.is_placeholder]
            if placeholders:
                summary_parts.append(f"  占位符数量: {len(placeholders)}")
        
        summary = "\n".join(summary_parts)
        logger.debug(f"--- [PPTParser]: Extracted text summary:\n{summary}")
        return summary
    
    def get_placeholder_mapping(self) -> Dict[int, List[Dict[str, Any]]]:
        """
        获取每张幻灯片的占位符映射
        
        Returns:
            字典，键是幻灯片索引，值是占位符信息列表
        """
        mapping = {}
        
        for idx, slide in enumerate(self.prs.slides):
            placeholders = []
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholder_info = {
                        "placeholder_id": shape.placeholder_format.idx,
                        "placeholder_type": str(shape.placeholder_format.type),
                        "has_text": False,
                        "text": ""
                    }
                    
                    if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                        placeholder_info["has_text"] = True
                        placeholder_info["text"] = shape.text_frame.text.strip()
                    
                    placeholders.append(placeholder_info)
            
            if placeholders:
                mapping[idx] = placeholders
        
        return mapping

