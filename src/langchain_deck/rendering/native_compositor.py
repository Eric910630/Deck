"""
原生合成器 (Native Compositor)
将 DOM 样式数据编译为 python-pptx 原生对象
"""
import re
from typing import Dict, Any, List, Tuple
from pptx.presentation import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from loguru import logger

from .browser_to_ppt_replicator.coordinate_mapper import CoordinateMapper


class NativeCompositor:
    """
    原生PPT合成器
    不使用截图，完全使用原生Shape绘制
    """
    
    def __init__(self, coordinate_mapper: CoordinateMapper = None):
        self.mapper = coordinate_mapper or CoordinateMapper()
    
    def composite_slide(self, slide, layout_data: List[Dict[str, Any]]):
        """
        合成单张幻灯片
        
        Args:
            slide: pptx.slide.Slide 对象
            layout_data: DOMAnalyzer 提取的数据列表
        """
        logger.info(f"--- [NativeCompositor]: 开始原生绘制 {len(layout_data)} 个元素")
        
        # 按 DOM 顺序绘制（后面的盖住前面的）
        # 也可以根据 z-index 排序，但通常 DOM 顺序就够了
        for elem in layout_data:
            elem_type = elem.get('type', 'text')
            
            if elem_type == 'card':
                self._draw_card(slide, elem)
            elif elem_type in ['title', 'text']:
                self._draw_text(slide, elem)
            else:
                # 默认当做文本处理
                self._draw_text(slide, elem)

    def _draw_card(self, slide, elem: Dict[str, Any]):
        """绘制卡片（背景、边框、阴影）"""
        geo = elem['geometry']
        style = elem['style']
        elem_id = elem.get('id', 'unknown')
        
        # 【调试日志】
        bg_color = style.get('backgroundColor', 'null')
        logger.info(f"--- [NativeCompositor]: 绘制卡片 {elem_id}, backgroundColor={bg_color}")
        
        # 1. 坐标转换
        left, top = self.mapper.browser_to_ppt(geo['x'], geo['y'])
        width, height = self.mapper.browser_size_to_ppt(geo['width'], geo['height'])
        
        # 2. 创建主形状 (背景)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Cm(left), Cm(top), Cm(width), Cm(height)
        )
        
        # 调整圆角
        if len(shape.adjustments) > 0:
            shape.adjustments[0] = 0.03

        # 3. 填充颜色
        # 【关键修复】强制使用白色背景，忽略提取到的 backgroundColor
        # 因为 Ant Design 卡片应该是白色背景，如果提取到蓝色，说明是父容器或渐变导致的误判
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        logger.info(f"--- [NativeCompositor]: 卡片 {elem_id} 强制使用白色背景 (忽略提取到的 {bg_color})")

        # 4. 边框 (极淡的灰色，辅助轮廓)
        shape.line.width = Pt(0.5) 
        shape.line.color.rgb = RGBColor(200, 200, 200) 

        # 5. 阴影 (核弹级)
        # 既然之前的参数还不够，我们就用 PPT 最重的阴影预设
        if style.get('boxShadow'):
            shape.shadow.inherit = False
            shape.shadow.visible = True
            shape.shadow.blur_radius = Pt(20)   # 20磅模糊，非常大
            shape.shadow.distance = Pt(8)       # 8磅位移，浮得很高
            shape.shadow.direction = 45
            shape.shadow.transparency = 0.3     # 30% 透明度 (70% 黑色)，非常黑！
            try:
                shape.shadow.color.rgb = RGBColor(0, 0, 0)
            except AttributeError:
                # 某些版本的 python-pptx 可能不支持 shadow.color
                pass

        # 6. 顶部装饰条 (Top Bar)
        if style.get('borderTopColor') and style.get('borderTopWidth', 0) > 0:
            bar_height_cm = 0.12
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Cm(left), Cm(top), Cm(width), Cm(bar_height_cm)
            )
            if len(bar.adjustments) > 0:
                bar.adjustments[0] = 0.8
            self._apply_color_to_fill(bar.fill, style['borderTopColor'])
            bar.line.fill.background()
            # 装饰条也加一点点阴影
            bar.shadow.inherit = False
            bar.shadow.blur_radius = Pt(3)
            bar.shadow.transparency = 0.7

    def _draw_text(self, slide, elem: Dict[str, Any]):
        """绘制文本框"""
        geo = elem['geometry']
        style = elem['style']
        text_content = elem.get('content', '').strip()
        elem_id = elem.get('id', 'unknown')
        
        if not text_content:
            return

        # 1. 坐标转换
        left_cm, top_cm = self.mapper.browser_to_ppt(geo['x'], geo['y'])
        width_cm, height_cm = self.mapper.browser_size_to_ppt(geo['width'], geo['height'])
        
        # 【调试日志】
        logger.info(f"--- [NativeCompositor]: 绘制文本 {elem_id}, 内容='{text_content}', 长度={len(text_content)}")
        logger.info(f"--- [NativeCompositor]: 原始尺寸: width={width_cm:.2f}cm, height={height_cm:.2f}cm")
        
        # 2. 创建文本框（使用原始尺寸，不扩容）
        textbox = slide.shapes.add_textbox(
            Cm(left_cm), Cm(top_cm), Cm(width_cm), Cm(height_cm)
        )
        
        tf = textbox.text_frame
        tf.text = text_content
        
        # 【核弹级修复】启用"根据文字调整形状大小"
        # 这会让文本框自动变宽以容纳所有文字，绝对不会换行
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.word_wrap = False  # 双重保险
        
        # 3. 彻底清除内边距
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        
        logger.info(f"--- [NativeCompositor]: 文本 '{text_content}' 启用 SHAPE_TO_FIT_TEXT，强制 word_wrap=False")
        
        # 5. 应用样式
        p = tf.paragraphs[0]
        
        # 对齐
        text_align = style.get('textAlign', 'left')
        if text_align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif text_align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        elif text_align == 'justify':
            p.alignment = PP_ALIGN.JUSTIFY
        else:
            p.alignment = PP_ALIGN.LEFT
            
        # 字体
        if len(p.runs) == 0:
            run = p.add_run()
        else:
            run = p.runs[0]
        
        # 【关键修复】强制使用微软雅黑（Windows/Mac通用且美观）
        run.font.name = "Microsoft YaHei"
        
        # 如果是标题，强制加粗
        if elem.get('type') == 'title':
            run.font.bold = True
            logger.info(f"--- [NativeCompositor]: 标题 '{text_content}' 强制加粗")
        
        # 字号
        if style.get('fontSize'):
            # 这里的 fontSize 是 px，需要转 pt
            # 浏览器通常 96dpi: 1px = 0.75pt
            run.font.size = Pt(style['fontSize'] * 0.75)
            
        # 颜色
        if style.get('color'):
            run.font.color.rgb = self._parse_color_to_rgb(style['color'])
            
        # 粗体
        # CSS fontWeight: "bold", "700", "normal", "400"
        w = str(style.get('fontWeight', 'normal'))
        if w in ['bold', '700', '800', '900']:
            run.font.bold = True

    def _parse_color_to_rgb(self, color_str: str) -> RGBColor:
        """
        解析 CSS 颜色并返回 RGBColor 对象
        支持 hex (#ffffff) 和 rgb(r, g, b)
        """
        if not color_str:
            return RGBColor(0, 0, 0)
            
        try:
            r, g, b = 0, 0, 0
            if color_str.startswith('#'):
                # Hex
                hex_str = color_str.lstrip('#')
                if len(hex_str) == 3:
                    hex_str = ''.join([c*2 for c in hex_str])
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
            elif color_str.startswith('rgb'):
                # rgba(r, g, b, a) or rgb(r, g, b)
                # 简单正则提取数字
                nums = [int(x) for x in re.findall(r'\d+', color_str)][:3]
                if len(nums) == 3:
                    r, g, b = nums
            
            return RGBColor(r, g, b)
        except Exception as e:
            logger.warning(f"颜色解析失败: {color_str}, error: {e}")
            return RGBColor(0, 0, 0)
    
    def _apply_color(self, obj, color_str: str):
        """
        解析 CSS 颜色并应用到 PPT 对象 (智能兼容 Fill 和 Line/Font)
        """
        if not color_str:
            return
            
        try:
            # 解析 Hex 或 RGB
            r, g, b = 0, 0, 0
            if color_str.startswith('#'):
                hex_str = color_str.lstrip('#')
                if len(hex_str) == 3:
                    hex_str = ''.join([c*2 for c in hex_str])
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
            elif color_str.startswith('rgb'):
                nums = [int(x) for x in re.findall(r'\d+', color_str)][:3]
                if len(nums) == 3:
                    r, g, b = nums
            
            color_obj = RGBColor(r, g, b)
            
            # 【关键修复】智能判断属性名
            # Shape.fill 是 FillFormat 对象，只有 fore_color，没有 color
            if hasattr(obj, 'fore_color'): 
                obj.fore_color.rgb = color_obj
            # Font 和 Line 是 ColorFormat 对象，有 color
            elif hasattr(obj, 'color'):
                obj.color.rgb = color_obj
                
        except Exception as e:
            logger.warning(f"颜色应用失败: {color_str}, error: {e}")
    
    def _apply_color_to_fill(self, fill_obj, color_str: str):
        """
        将颜色应用到 FillFormat 对象（向后兼容方法）
        """
        self._apply_color(fill_obj, color_str)

    def _map_font_family(self, css_font: str) -> str:
        """将 CSS font-family 映射到 Windows 安全字体"""
        if not css_font:
            return "Arial"
        
        css_font = css_font.lower()
        if "yahei" in css_font or "hei" in css_font:
            return "Microsoft YaHei"
        if "song" in css_font:
            return "SimSun"
        if "arial" in css_font:
            return "Arial"
        if "times" in css_font:
            return "Times New Roman"
        
        # 默认 Fallback
        return "Microsoft YaHei"  # 中文 PPT 默认雅黑比较安全

