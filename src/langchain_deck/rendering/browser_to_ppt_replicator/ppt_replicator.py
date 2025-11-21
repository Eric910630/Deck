"""
PPT复刻器
将浏览器渲染结果复刻到PPT
"""

from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from loguru import logger

from .coordinate_mapper import CoordinateMapper


class PPTReplicator:
    """
    PPT复刻器
    将容器图片和文本插入到PPT的相同位置
    """
    
    def __init__(self, coordinate_mapper: CoordinateMapper, output_path: Optional[Path] = None, prs: Optional[Presentation] = None):
        """
        初始化PPT复刻器
        
        Args:
            coordinate_mapper: 坐标映射器
            output_path: 输出PPT路径
            prs: 现有的Presentation对象（如果提供，使用它而不是创建新的）
        """
        self.mapper = coordinate_mapper
        
        if prs is not None:
            # 使用现有的Presentation对象（用于多张幻灯片）
            self.prs = prs
        else:
            # 创建新的Presentation对象
            self.prs = Presentation()
            # 设置16:9尺寸
            self.prs.slide_width = Cm(coordinate_mapper.PPT_WIDTH_CM)
            self.prs.slide_height = Cm(coordinate_mapper.PPT_HEIGHT_CM)
        
        self.output_path = output_path
        logger.info(f"--- [PPTReplicator]: Initialized (size: {coordinate_mapper.PPT_WIDTH_CM}cm x {coordinate_mapper.PPT_HEIGHT_CM}cm)")
    
    def replicate_slide(
        self,
        containers: List[Dict[str, Any]],
        texts: List[Dict[str, Any]],
        use_hybrid_rendering: bool = True
    ) -> None:
        """
        复刻一张幻灯片
        
        【新架构原则】：混合渲染法
        - 容器用图片（浏览器截图，保留完美样式）
        - 内容用原生文本（保证可编辑性）
        
        Args:
            containers: 容器列表（包含图片路径和位置）
            texts: 文本列表（包含内容和样式）
            use_hybrid_rendering: 是否使用混合渲染（默认True）
        """
        logger.info(f"--- [PPTReplicator]: Replicating slide with {len(containers)} containers and {len(texts)} texts")
        logger.info(f"--- [PPTReplicator]: 使用混合渲染: {use_hybrid_rendering}")
        
        # 创建空白幻灯片
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 空白布局
        
        if use_hybrid_rendering:
            # 【混合渲染法】：容器用图片，文本用原生
            # 1. 插入容器图片（底层，按z-index从后往前）
            sorted_containers = sorted(containers, key=lambda c: c.get('z_index', 0))
            logger.info(f"--- [PPTReplicator]: 【混合渲染】插入{len(sorted_containers)}个容器图片（作为背景）")
            for container in sorted_containers:
                self._insert_container_image(slide, container)
            
            # 2. 插入文本（顶层，可编辑）
            # 对于卡片内的文本，需要根据容器位置计算文本位置
            logger.info(f"--- [PPTReplicator]: 【混合渲染】插入{len(texts)}个文本元素（可编辑）")
            for text in texts:
                # 检查文本是否在容器内
                text_in_container = False
                container_for_text = None
                
                text_x = text['position']['x']
                text_y = text['position']['y']
                text_w = text['size']['width']
                text_h = text['size']['height']
                
                for container in sorted_containers:
                    cont_x = container['position']['x']
                    cont_y = container['position']['y']
                    cont_w = container['size']['width']
                    cont_h = container['size']['height']
                    
                    # 检查文本是否在容器内（允许小范围重叠）
                    if (text_x >= cont_x - 10 and text_y >= cont_y - 10 and
                        text_x + text_w <= cont_x + cont_w + 10 and
                        text_y + text_h <= cont_y + cont_h + 10):
                        text_in_container = True
                        container_for_text = container
                        break
                
                if text_in_container and container_for_text:
                    # 文本在容器内：在容器上方插入文本（可编辑）
                    # 注意：容器图片已经包含了样式，文本只需要内容
                    self._insert_text_on_container(slide, text, container_for_text)
                else:
                    # 文本不在容器内：独立插入
                    self._insert_text(slide, text)
        else:
            # 【旧方法】：只插入容器图片，文本已包含在图片中
            sorted_containers = sorted(containers, key=lambda c: c.get('z_index', 0))
            logger.info(f"--- [PPTReplicator]: 插入{len(sorted_containers)}个容器图片（文本已包含在容器图片中）")
            for container in sorted_containers:
                self._insert_container_image(slide, container)
            
            # 只插入不在容器内的文本
            texts_to_insert = []
            for text in texts:
                text_in_container = False
                text_x = text['position']['x']
                text_y = text['position']['y']
                text_w = text['size']['width']
                text_h = text['size']['height']
                
                for container in sorted_containers:
                    cont_x = container['position']['x']
                    cont_y = container['position']['y']
                    cont_w = container['size']['width']
                    cont_h = container['size']['height']
                    
                    if (text_x >= cont_x - 10 and text_y >= cont_y - 10 and
                        text_x + text_w <= cont_x + cont_w + 10 and
                        text_y + text_h <= cont_y + cont_h + 10):
                        text_in_container = True
                        break
                
                if not text_in_container:
                    texts_to_insert.append(text)
            
            if texts_to_insert:
                logger.info(f"--- [PPTReplicator]: 插入{len(texts_to_insert)}个独立文本元素")
                for text in texts_to_insert:
                    self._insert_text(slide, text)
        
        logger.info("--- [PPTReplicator]: Slide replicated")
    
    def _insert_container_image(self, slide, container: Dict[str, Any]):
        """插入容器图片"""
        try:
            image_path = container['image_path']
            if not Path(image_path).exists():
                logger.warning(f"--- [PPTReplicator]: Image not found: {image_path}")
                return
            
            # 转换坐标
            browser_x = container['position']['x']
            browser_y = container['position']['y']
            browser_w = container['size']['width']
            browser_h = container['size']['height']
            
            ppt_x, ppt_y = self.mapper.browser_to_ppt(browser_x, browser_y)
            ppt_width, ppt_height = self.mapper.browser_size_to_ppt(browser_w, browser_h)
            
            # 【日志探针】记录插入信息
            logger.info(f"--- [PPTReplicator]: 【容器图片】")
            logger.info(f"    浏览器位置: ({browser_x:.1f}px, {browser_y:.1f}px)")
            logger.info(f"    浏览器尺寸: {browser_w:.1f}px × {browser_h:.1f}px")
            logger.info(f"    PPT位置: ({ppt_x:.2f}cm, {ppt_y:.2f}cm)")
            logger.info(f"    PPT尺寸: {ppt_width:.2f}cm × {ppt_height:.2f}cm")
            logger.info(f"    覆盖区域: x=[{ppt_x:.2f}, {ppt_x+ppt_width:.2f}], y=[{ppt_y:.2f}, {ppt_y+ppt_height:.2f}]")
            
            # 检查是否超出边界
            if ppt_x < 0 or ppt_y < 0:
                logger.warning(f"--- [PPTReplicator]: ⚠️ 位置超出边界 (x={ppt_x:.2f}, y={ppt_y:.2f})")
            if ppt_x + ppt_width > self.mapper.PPT_WIDTH_CM or ppt_y + ppt_height > self.mapper.PPT_HEIGHT_CM:
                logger.warning(f"--- [PPTReplicator]: ⚠️ 尺寸超出边界")
            
            # 插入图片
            slide.shapes.add_picture(
                image_path,
                Cm(ppt_x),
                Cm(ppt_y),
                Cm(ppt_width),
                Cm(ppt_height)
            )
            
            logger.debug(f"--- [PPTReplicator]: Inserted container image at ({ppt_x:.2f}cm, {ppt_y:.2f}cm)")
        except Exception as e:
            logger.warning(f"--- [PPTReplicator]: Failed to insert container image: {e}")
            import traceback
            logger.debug(traceback.format_exc())
    
    def _insert_text_on_container(
        self, 
        slide, 
        text: Dict[str, Any], 
        container: Dict[str, Any]
    ):
        """
        在容器上方插入文本（混合渲染法）
        
        文本位置相对于容器，需要考虑容器的内边距
        
        Args:
            slide: PPT 幻灯片对象
            text: 文本信息
            container: 容器信息
        """
        try:
            # 计算文本相对于容器的位置
            # 假设容器有 24px 内边距（padding-lg）
            container_padding = 24  # px
            
            # 文本在浏览器中的绝对位置
            text_x = text['position']['x']
            text_y = text['position']['y']
            text_w = text['size']['width']
            text_h = text['size']['height']
            
            # 容器在浏览器中的绝对位置
            cont_x = container['position']['x']
            cont_y = container['position']['y']
            
            # 转换坐标
            ppt_x, ppt_y = self.mapper.browser_to_ppt(text_x, text_y)
            ppt_width, ppt_height = self.mapper.browser_size_to_ppt(text_w, text_h)
            
            logger.info(f"--- [PPTReplicator]: 【混合渲染】在容器上方插入文本")
            logger.info(f"    文本浏览器位置: ({text_x:.1f}px, {text_y:.1f}px)")
            logger.info(f"    容器浏览器位置: ({cont_x:.1f}px, {cont_y:.1f}px)")
            logger.info(f"    PPT位置: ({ppt_x:.2f}cm, {ppt_y:.2f}cm)")
            
            # 插入文本框（可编辑）
            textbox = slide.shapes.add_textbox(
                Cm(ppt_x),
                Cm(ppt_y),
                Cm(ppt_width),
                Cm(ppt_height)
            )
            
            # 设置文本内容
            tf = textbox.text_frame
            tf.text = text.get('content', '')
            tf.word_wrap = True
            
            # 设置文本样式
            if tf.paragraphs:
                para = tf.paragraphs[0]
                text_style = text.get('style', {})
                
                # 对齐方式
                text_align = text_style.get('textAlign', 'left')
                if text_align == 'center':
                    para.alignment = PP_ALIGN.CENTER
                elif text_align == 'right':
                    para.alignment = PP_ALIGN.RIGHT
                else:
                    para.alignment = PP_ALIGN.LEFT
                
                if para.runs:
                    run = para.runs[0]
                    # 字号（从px转换为pt）
                    font_size_str = text_style.get('fontSize', '14px')
                    font_size_pt = self._px_to_pt(font_size_str)
                    run.font.size = Pt(font_size_pt)
                    
                    # 字重
                    font_weight = text_style.get('fontWeight', '400')
                    run.font.bold = (font_weight == '600' or font_weight == '700' or font_weight == 'bold')
                    
                    # 颜色
                    color_str = text_style.get('color', '#000000')
                    run.font.color.rgb = self._parse_color(color_str)
            
            logger.debug(f"--- [PPTReplicator]: Inserted text on container at ({ppt_x:.2f}cm, {ppt_y:.2f}cm)")
        except Exception as e:
            logger.warning(f"--- [PPTReplicator]: Failed to insert text on container: {e}")
            import traceback
            logger.debug(traceback.format_exc())
    
    def _insert_text(self, slide, text: Dict[str, Any]):
        """插入文本"""
        try:
            browser_x = text['position']['x']
            browser_y = text['position']['y']
            browser_w = text['size']['width']
            browser_h = text['size']['height']
            
            # 转换坐标
            ppt_x, ppt_y = self.mapper.browser_to_ppt(browser_x, browser_y)
            ppt_width, ppt_height = self.mapper.browser_size_to_ppt(browser_w, browser_h)
            
            # 【日志探针】记录插入信息
            logger.info(f"--- [PPTReplicator]: 【文本】")
            logger.info(f"    浏览器位置: ({browser_x:.1f}px, {browser_y:.1f}px)")
            logger.info(f"    浏览器尺寸: {browser_w:.1f}px × {browser_h:.1f}px")
            logger.info(f"    PPT位置: ({ppt_x:.2f}cm, {ppt_y:.2f}cm)")
            logger.info(f"    PPT尺寸: {ppt_width:.2f}cm × {ppt_height:.2f}cm")
            logger.info(f"    文本内容: {text['text'][:50]}...")
            logger.info(f"    覆盖区域: x=[{ppt_x:.2f}, {ppt_x+ppt_width:.2f}], y=[{ppt_y:.2f}, {ppt_y+ppt_height:.2f}]")
            
            # 检查是否超出边界
            if ppt_x < 0 or ppt_y < 0:
                logger.warning(f"--- [PPTReplicator]: ⚠️ 位置超出边界 (x={ppt_x:.2f}, y={ppt_y:.2f})")
            if ppt_x + ppt_width > self.mapper.PPT_WIDTH_CM or ppt_y + ppt_height > self.mapper.PPT_HEIGHT_CM:
                logger.warning(f"--- [PPTReplicator]: ⚠️ 尺寸超出边界")
            
            # 创建文本框
            textbox = slide.shapes.add_textbox(
                Cm(ppt_x),
                Cm(ppt_y),
                Cm(ppt_width),
                Cm(ppt_height)
            )
            
            # 设置文本
            text_frame = textbox.text_frame
            text_frame.text = text['text']
            text_frame.word_wrap = True
            
            # 应用样式
            style = text['style']
            para = text_frame.paragraphs[0]
            run = para.runs[0]
            
            # 字体
            run.font.name = style.get('font_family', 'Arial')
            run.font.size = Pt(style.get('font_size_pt', 14))
            run.font.bold = style.get('is_bold', False)
            
            # 颜色
            color_hex = style.get('color_hex', '#000000')
            run.font.color.rgb = self._hex_to_rgb(color_hex)
            
            # 对齐方式
            text_align = style.get('text_align', 'left')
            if text_align == 'center':
                para.alignment = PP_ALIGN.CENTER
            elif text_align == 'right':
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.LEFT
            
            logger.debug(f"--- [PPTReplicator]: Inserted text at ({ppt_x:.2f}cm, {ppt_y:.2f}cm): {text['text'][:30]}...")
        except Exception as e:
            logger.warning(f"--- [PPTReplicator]: Failed to insert text: {e}")
            import traceback
            logger.debug(traceback.format_exc())
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将hex颜色转换为RGBColor"""
        try:
            if hex_color.startswith('#'):
                r = int(hex_color[1:3], 16)
                g = int(hex_color[3:5], 16)
                b = int(hex_color[5:7], 16)
                return RGBColor(r, g, b)
        except:
            pass
        return RGBColor(0, 0, 0)  # 默认黑色
    
    def save(self, output_path: Optional[Path] = None) -> Path:
        """
        保存PPT
        
        Args:
            output_path: 输出路径（如果未提供，使用初始化时的路径）
            
        Returns:
            保存的文件路径
        """
        if output_path is None:
            output_path = self.output_path
        
        if output_path is None:
            from datetime import datetime
            output_path = Path(f"replicated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
        
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(output_path))
        logger.info(f"--- [PPTReplicator]: PPT saved to {output_path}")
        return output_path

