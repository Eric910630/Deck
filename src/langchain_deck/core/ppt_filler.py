"""
PPT内容填充器
根据PPT框架和LLM生成的内容，填充完整的PPT
"""

from pathlib import Path
from typing import Dict, List, Any, Optional
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from loguru import logger

from ..analysis.ppt_parser import PPTParser
from .llm_service import LLMService, create_llm_service
from ..analysis.enhanced_ppt_parser import EnhancedPPTParser
from ..analysis.human_centered_analyzer import HumanCenteredAnalyzer
from ..planning.content_strategy_generator import ContentStrategyGenerator
from ..planning.content_polisher import ContentPolisher
from ..planning.presentation_planner import PresentationPlanner
from ..planning.layout_planner import LayoutPlanner
from ..rendering.html_generator import HTMLGenerator
from ..theme.ant_design_theme import ant_design_theme
from ..rendering.browser_to_ppt_replicator import BrowserToPPTReplicator

# 【新增】导入 Native Compiler 组件
from ..rendering.dom_analyzer import DOMAnalyzer
from ..rendering.native_compositor import NativeCompositor
from ..rendering.browser_to_ppt_replicator.browser_renderer import BrowserRenderer
from ..rendering.browser_to_ppt_replicator.coordinate_mapper import CoordinateMapper


class PPTFiller:
    """
    PPT内容填充器
    根据框架PPT和用户需求，使用LLM生成内容并填充
    """
    
    def __init__(
        self,
        framework_path: str,
        llm_service: Optional[LLMService] = None,
        use_browser_rendering: bool = False
    ):
        """
        初始化PPT填充器
        
        Args:
            framework_path: 框架PPT文件路径
            llm_service: LLM服务实例
            use_browser_rendering: 是否使用浏览器渲染（Ant Design规范）
        """
        self.parser = PPTParser(framework_path)
        self.framework_path = Path(framework_path)
        self.use_browser_rendering = use_browser_rendering
        
        if llm_service is None:
            self.llm_service = create_llm_service(use_async=True)
            if self.llm_service is None:
                logger.warning("--- [PPTFiller]: LLM service not available")
        else:
            self.llm_service = llm_service
        
        # 如果使用浏览器渲染，初始化相关组件
        if use_browser_rendering:
            self.html_generator = HTMLGenerator()
            self.replicator = BrowserToPPTReplicator()
            logger.info("--- [PPTFiller]: Initialized with browser rendering enabled")
        else:
            self.html_generator = None
            self.replicator = None
            logger.info("--- [PPTFiller]: Initialized")
    
    async def fill_from_prompt(
        self,
        prompt: str,
        output_path: Optional[str] = None,
        preserve_structure: bool = True,
        use_enhanced_analysis: bool = True,
        use_browser_rendering: Optional[bool] = None,
        skip_ppt_conversion: bool = False
    ) -> str:
        """
        根据用户提示填充PPT内容
        
        Args:
            prompt: 用户需求描述
            output_path: 输出文件路径，如果为None则自动生成
            preserve_structure: 是否保持原有结构
            use_enhanced_analysis: 是否使用增强分析（人类中心化分析）
            use_browser_rendering: 是否使用浏览器渲染（如果为None，使用初始化时的设置）
            skip_ppt_conversion: 是否跳过HTML到PPT的转换（仅生成HTML文件）
            
        Returns:
            生成的PPT文件路径（如果skip_ppt_conversion=True，返回HTML目录路径）
        """
        if not self.llm_service:
            raise ValueError("LLM service is required for content generation")
        
        # 确定是否使用浏览器渲染
        if use_browser_rendering is None:
            use_browser_rendering = self.use_browser_rendering
        
        # 如果使用浏览器渲染，使用新的流程
        if use_browser_rendering:
            if not self.html_generator or not self.replicator:
                raise ValueError("Browser rendering components not initialized. Set use_browser_rendering=True in __init__")
            return await self._fill_with_browser_rendering(
                prompt, output_path, use_enhanced_analysis, skip_ppt_conversion
            )
        
        if use_enhanced_analysis:
            # 使用增强分析流程
            logger.info("--- [PPTFiller]: Using enhanced analysis workflow...")
            
            # 1. 提取增强结构
            logger.info("--- [PPTFiller]: Extracting enhanced structure...")
            enhanced_parser = EnhancedPPTParser(str(self.framework_path))
            enhanced_structure = enhanced_parser.extract_structure_enhanced()
            
            # 2. 人类中心化分析
            logger.info("--- [PPTFiller]: Performing human-centered analysis...")
            analyzer = HumanCenteredAnalyzer(enhanced_structure)
            human_analysis = analyzer.analyze_all()
            
            # 3. 生成内容策略
            logger.info("--- [PPTFiller]: Generating content strategy...")
            strategy_gen = ContentStrategyGenerator(human_analysis)
            content_strategy = strategy_gen.generate_strategy()
            
            # 4. 使用策略生成内容
            logger.info("--- [PPTFiller]: Generating content by sections...")
            content_map = await self._generate_content_by_sections(
                human_analysis=human_analysis,
                content_strategy=content_strategy,
                user_prompt=prompt
            )
        else:
            # 使用原有流程（向后兼容）
            logger.info("--- [PPTFiller]: Using legacy workflow...")
            structure = self.parser.extract_structure()
            text_summary = self.parser.extract_text_summary()
            placeholder_mapping = self.parser.get_placeholder_mapping()
            
            content_map = await self._generate_content_for_framework(
                prompt=prompt,
                structure=structure,
                text_summary=text_summary,
                placeholder_mapping=placeholder_mapping
            )
        
        # 填充PPT
        logger.info("--- [PPTFiller]: Filling PPT with generated content...")
        output_path = output_path or self._generate_output_path()
        
        # 传递分析结果和策略（如果使用了增强分析）
        if use_enhanced_analysis:
            self._fill_ppt(content_map, output_path, preserve_structure, human_analysis, content_strategy)
        else:
            self._fill_ppt(content_map, output_path, preserve_structure)
        
        logger.success(f"--- [PPTFiller]: PPT filled and saved to {output_path}")
        return str(output_path)
    
    async def _fill_with_browser_rendering(
        self,
        prompt: str,
        output_path: Optional[str] = None,
        use_enhanced_analysis: bool = True,
        skip_ppt_conversion: bool = False
    ) -> str:
        """
        使用浏览器渲染方式填充PPT（Ant Design规范）
        
        Args:
            prompt: 用户提示
            output_path: 输出路径
            use_enhanced_analysis: 是否使用增强分析
            skip_ppt_conversion: 是否跳过HTML到PPT的转换（仅生成HTML文件）
            
        Returns:
            输出文件路径（如果skip_ppt_conversion=True，返回HTML目录路径）
        """
        logger.info("="*80)
        logger.info("--- [PPTFiller]: 使用浏览器渲染方式（Ant Design规范）")
        logger.info("="*80)
        
        # 1. 提取框架结构
        if use_enhanced_analysis:
            logger.info("--- [PPTFiller]: 提取增强结构...")
            enhanced_parser = EnhancedPPTParser(str(self.framework_path))
            enhanced_structure = enhanced_parser.extract_structure_enhanced()
            
            # 2. 人类中心化分析
            # 如果prompt包含docx内容，使用docx内容进行分析，否则使用框架PPT结构
            logger.info("--- [PPTFiller]: 执行人类中心化分析...")
            
            # 检查prompt是否包含docx内容（通过检查是否包含"【文档内容】"标记）
            if "【文档内容】" in prompt or "文档内容" in prompt:
                # 从prompt中提取docx内容
                import re
                docx_content_match = re.search(r'【文档内容】\s*\n(.*?)(?=\n【|$)', prompt, re.DOTALL)
                if docx_content_match:
                    docx_content = docx_content_match.group(1).strip()
                    logger.info(f"--- [PPTFiller]: 检测到docx内容，长度: {len(docx_content)}字符")
                    logger.info(f"--- [PPTFiller]: 使用docx内容进行人类中心化分析（而不是框架PPT内容）")
                    # 创建基于docx内容的结构数据
                    docx_structure = self._create_structure_from_docx_content(docx_content, enhanced_structure)
                    analyzer = HumanCenteredAnalyzer(docx_structure, raw_text=docx_content, llm_service=self.llm_service)
                else:
                    logger.warning("--- [PPTFiller]: 未找到docx内容，使用框架PPT结构进行分析")
                    analyzer = HumanCenteredAnalyzer(enhanced_structure, llm_service=self.llm_service)
            else:
                logger.info("--- [PPTFiller]: 使用框架PPT结构进行分析")
                analyzer = HumanCenteredAnalyzer(enhanced_structure, llm_service=self.llm_service)
            
            human_analysis = await analyzer.analyze_all()
            
            # 3. 生成内容策略
            logger.info("--- [PPTFiller]: 生成内容策略...")
            strategy_generator = ContentStrategyGenerator(human_analysis)
            content_strategy = strategy_generator.generate_strategy()
            
            # 4. 智能识别支撑材料
            logger.info("--- [PPTFiller]: 智能识别支撑材料...")
            from ..analysis.supporting_materials_analyzer import SupportingMaterialsAnalyzer
            materials_analyzer = SupportingMaterialsAnalyzer(self.llm_service)
            
            supporting_materials = human_analysis.get("layer_4_supporting_materials", {}).get("data", {})
            raw_data_points = supporting_materials.get("materials", {}).get("data_points", [])
            raw_cases = supporting_materials.get("materials", {}).get("cases", [])
            
            # 智能识别数据点和案例
            intelligent_data_points = await materials_analyzer.intelligently_identify_data_points(raw_data_points)
            intelligent_cases = await materials_analyzer.intelligently_identify_cases(raw_cases)
            
            logger.info(f"--- [PPTFiller]: 识别出{len(intelligent_data_points)}个数据点，{len(intelligent_cases)}个案例")
            
            # 5. 逐板块生成内容（整合支撑材料）
            logger.info("--- [PPTFiller]: 逐板块生成内容（整合支撑材料）...")
            generation_result = await self._generate_content_by_sections(
                human_analysis, content_strategy, prompt,
                intelligent_data_points, intelligent_cases
            )
            
            # 【新增】提取内容映射、润色结果、布局规划结果和颜色配置结果
            if isinstance(generation_result, dict):
                content_map = generation_result.get('content_map', {})
                polished_slides = generation_result.get('polished_slides', [])
                presentation_plans = generation_result.get('presentation_plans', [])
                layout_plans = generation_result.get('layout_plans', [])
                color_configs = generation_result.get('color_configs', [])
            else:
                # 向后兼容：如果返回的是旧格式（只有content_map）
                content_map = generation_result
                polished_slides = []
                presentation_plans = []
                layout_plans = []
                color_configs = []
            
            # 【探针】检查内容映射
            logger.info("="*80)
            logger.info("--- [PPTFiller]: 【探针】内容映射总览")
            logger.info("="*80)
            logger.info(f"   总内容映射项数: {len(content_map)}")
            logger.info(f"   润色幻灯片数量: {len(polished_slides)}")
            logger.info(f"   展示策划数量: {len(presentation_plans)}")
            logger.info(f"   布局规划数量: {len(layout_plans)}")
            logger.info(f"   颜色配置数量: {len(color_configs)}")
            slides_in_map = set()
            for key in content_map.keys():
                if 'slide_' in key:
                    try:
                        slide_idx = int(key.split('_')[1])
                        slides_in_map.add(slide_idx)
                    except:
                        pass
            logger.info(f"   涉及幻灯片: {sorted(slides_in_map)}")
            logger.info(f"   内容映射键列表:")
            for key in sorted(content_map.keys()):
                logger.info(f"     - {key}: {len(content_map[key])}字符")
            logger.info("="*80)
        else:
            # 使用简单方式生成内容
            logger.info("--- [PPTFiller]: 使用简单方式生成内容...")
            structure = self.parser.extract_structure()
            text_summary = self.parser.extract_text_summary()
            placeholder_mapping = self.parser.get_placeholder_mapping()
            content_map = await self._generate_content_for_framework(
                prompt=prompt,
                structure=structure,
                text_summary=text_summary,
                placeholder_mapping=placeholder_mapping
            )
        
        # 5. 生成HTML（Ant Design规范）
        # 【新增】优先使用布局规划结果和颜色配置生成HTML，如果没有则使用content_map
        output_path_obj = Path(output_path) if output_path else Path(self._generate_output_path())
        html_output_dir = output_path_obj.parent / "html_output"
        html_output_dir.mkdir(parents=True, exist_ok=True)
        
        if layout_plans and polished_slides:
            logger.info("--- [PPTFiller]: 使用布局规划和颜色配置生成HTML（精确布局+颜色）...")
            
            # 【新增】生成合并的HTML文件
            merged_html = self.html_generator.generate_merged_html(
                layout_plans=layout_plans,
                polished_slides=polished_slides,
                color_configs=color_configs if color_configs else None
            )
            
            # 保存合并的HTML文件
            merged_html_file = html_output_dir / "presentation.html"
            merged_html_file.write_text(merged_html, encoding='utf-8')
            logger.info(f"--- [PPTFiller]: ✅ 保存合并HTML到: {merged_html_file}")
            
            # 同时保存单独的HTML文件（用于调试）
            html_contents = self.html_generator.generate_from_layout_plan(
                layout_plans=layout_plans,
                polished_slides=polished_slides,
                color_configs=color_configs if color_configs else None
            )
            for idx, html_content in enumerate(html_contents):
                html_file = html_output_dir / f"slide_{idx:03d}.html"
                html_file.write_text(html_content, encoding='utf-8')
                logger.info(f"--- [PPTFiller]: ✅ 保存单独HTML到: {html_file}")
        else:
            logger.info("--- [PPTFiller]: 使用内容映射生成HTML（标准方式）...")
            html_contents = self.html_generator.generate_from_content_map(content_map)
            
            # 【修复】支持多张幻灯片：html_contents现在可能是列表
            if not isinstance(html_contents, list):
                html_contents = [html_contents]
            
            for idx, html_content in enumerate(html_contents):
                html_file = html_output_dir / f"slide_{idx:03d}.html"
                html_file.write_text(html_content, encoding='utf-8')
                logger.info(f"--- [PPTFiller]: ✅ 保存HTML到: {html_file}")
        
        logger.info(f"--- [PPTFiller]: 生成了{len(html_contents) if isinstance(html_contents, list) else 1}张HTML幻灯片")
        
        # 【新增】如果跳过PPT转换，直接返回HTML目录路径
        if skip_ppt_conversion:
            logger.info("="*80)
            logger.info("--- [PPTFiller]: ⏭️  跳过HTML到PPT的转换（仅生成HTML）")
            logger.info(f"--- [PPTFiller]: ✅ HTML文件已保存到: {html_output_dir}")
            logger.info(f"--- [PPTFiller]: 共生成 {len(html_contents)} 张HTML幻灯片")
            logger.info("="*80)
            logger.success(f"--- [PPTFiller]: HTML生成完成，保存到 {html_output_dir}")
            return str(html_output_dir)
        
        # ==========================================
        # 6. 【升级】Native Shape Compiler 流程
        # ==========================================
        logger.info("--- [PPTFiller]: 启动 Native Shape Compiler (原生重绘)...")
        output_path = output_path or self._generate_output_path()
        output_path_obj = Path(output_path)
        
        # 初始化 PPT (16:9)
        from pptx import Presentation
        from pptx.util import Cm
        prs = Presentation()
        prs.slide_width = Cm(33.867)
        prs.slide_height = Cm(19.05)
        
        # 初始化 Native Compiler 组件
        browser_renderer = BrowserRenderer()
        coordinate_mapper = CoordinateMapper()
        dom_analyzer = DOMAnalyzer()
        native_compositor = NativeCompositor(coordinate_mapper)
        
        try:
            for slide_idx, html_content in enumerate(html_contents):
                logger.info(f"--- [PPTFiller]: 编译第 {slide_idx + 1}/{len(html_contents)} 张幻灯片...")
                
                # 6.1 浏览器渲染 (Headless)
                page = await browser_renderer.render_html(html_content)
                
                try:
                    # 6.2 提取 DOM 样式数据 (不截图)
                    layout_data = await dom_analyzer.extract_layout_data(page)
                    logger.info(f"--- [PPTFiller]: 提取到 {len(layout_data)} 个原生元素数据")
                    
                    # 6.3 创建空白幻灯片
                    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
                    
                    # 6.4 原生绘制 (Native Draw)
                    native_compositor.composite_slide(slide, layout_data)
                    
                    logger.success(f"--- [PPTFiller]: 幻灯片 {slide_idx + 1} 原生编译完成")
                    
                finally:
                    await page.close()
            
            # 7. 保存 PPT
            prs.save(str(output_path_obj))
            replicated_path = output_path_obj
            logger.success(f"--- [PPTFiller]: 全局编译完成，文件已保存: {replicated_path}")
            logger.info(f"--- [PPTFiller]: 成功生成{len(prs.slides)}张PPT幻灯片")
            
        except Exception as e:
            logger.error(f"--- [PPTFiller]: Native Compiler 运行失败: {e}", exc_info=True)
            raise
        finally:
            await browser_renderer.close()
        
        # 7. 整合图表（如果有可可视化数据）
        if intelligent_data_points:
            logger.info("--- [PPTFiller]: 整合图表...")
            from ..charts.chart_integrator import ChartIntegrator
            chart_integrator = ChartIntegrator(
                materials_analyzer=materials_analyzer
            )
            
            # 打开PPT并插入图表
            from pptx import Presentation
            prs = Presentation(str(replicated_path))
            chart_count = await chart_integrator.integrate_charts(
                prs,
                intelligent_data_points,
                output_dir=replicated_path.parent / "charts"
            )
            prs.save(str(replicated_path))
            
            if chart_count > 0:
                logger.info(f"--- [PPTFiller]: 成功整合{chart_count}个图表")
        
        logger.success(f"--- [PPTFiller]: PPT filled and saved to {replicated_path}")
        return str(replicated_path)
    
    async def _generate_content_for_framework(
        self,
        prompt: str,
        structure: Dict[str, Any],
        text_summary: str,
        placeholder_mapping: Dict[int, List[Dict[str, Any]]]
    ) -> Dict[str, str]:
        """
        为框架生成内容映射
        
        Args:
            prompt: 用户需求
            structure: PPT结构信息
            text_summary: 文本摘要
            placeholder_mapping: 占位符映射
            
        Returns:
            内容映射字典，键是占位符标识，值是生成的内容
        """
        system_prompt = """你是一个专业的PPT内容创作助手。你的任务是根据用户需求和PPT框架结构，为每张幻灯片的占位符生成合适的内容。

要求：
1. 理解PPT框架的结构和现有内容（如果有）
2. 根据用户需求生成专业、相关的内容
3. 为每个占位符生成合适的内容
4. 保持内容的逻辑连贯性和专业性
5. 标题要简洁有力，正文要详细但不过长

输出格式（JSON）：
{
  "slide_0_placeholder_0": "标题内容",
  "slide_0_placeholder_1": "正文内容",
  "slide_1_placeholder_0": "标题内容",
  ...
}

占位符标识格式：slide_{幻灯片索引}_placeholder_{占位符ID}"""
        
        user_prompt = f"""PPT框架信息：
{text_summary}

用户需求：{prompt}

请为每张幻灯片的占位符生成合适的内容。如果占位符已有内容，可以基于现有内容进行扩展或优化。"""
        
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]
        
        try:
            response = await self.llm_service.chat_completion_async(
                messages=messages,
                temperature=0.7,
                max_tokens=4000
            )
            
            # 解析JSON响应
            import json
            import re
            
            # 尝试提取JSON
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if json_match:
                content_map = json.loads(json_match.group(0))
            else:
                # 如果无法解析，尝试手动构建
                logger.warning("--- [PPTFiller]: Failed to parse JSON, using fallback")
                content_map = self._fallback_content_generation(placeholder_mapping, prompt)
            
            return content_map
            
        except Exception as e:
            logger.error(f"--- [PPTFiller]: Failed to generate content: {e}", exc_info=True)
            # 使用fallback
            return self._fallback_content_generation(placeholder_mapping, prompt)
    
    def _fallback_content_generation(
        self,
        placeholder_mapping: Dict[int, List[Dict[str, Any]]],
        prompt: str
    ) -> Dict[str, str]:
        """
        Fallback内容生成（如果LLM失败）
        
        Args:
            placeholder_mapping: 占位符映射
            prompt: 用户需求
            
        Returns:
            内容映射字典
        """
        content_map = {}
        for slide_idx, placeholders in placeholder_mapping.items():
            for placeholder in placeholders:
                key = f"slide_{slide_idx}_placeholder_{placeholder['placeholder_id']}"
                if placeholder.get("has_text"):
                    # 如果有现有文本，保留
                    content_map[key] = placeholder["text"]
                else:
                    # 生成占位文本
                    content_map[key] = f"[需要填充: {prompt[:50]}...]"
        
        return content_map
    
    def _fill_ppt(
        self,
        content_map: Dict[str, str],
        output_path: str,
        preserve_structure: bool = True,
        human_analysis: Optional[Dict[str, Any]] = None,
        content_strategy: Optional[Dict[str, Any]] = None
    ):
        """
        填充PPT内容
        
        Args:
            content_map: 内容映射字典
            output_path: 输出路径
            preserve_structure: 是否保持结构
        """
        logger.info("="*80)
        logger.info("--- [PPTFiller]: 开始填充PPT流程")
        logger.info("="*80)
        
        # 复制框架PPT
        from shutil import copy
        output_path_obj = Path(output_path)
        output_path_obj.parent.mkdir(parents=True, exist_ok=True)
        copy(self.framework_path, output_path)
        logger.info(f"--- [PPTFiller]: 已复制框架PPT到: {output_path}")
        
        # 打开复制的PPT
        prs = Presentation(str(output_path))
        
        # 【日志探针1】检查原始PPT尺寸
        original_width_emu = prs.slide_width
        original_height_emu = prs.slide_height
        original_width_cm = float(original_width_emu) / 360000
        original_height_cm = float(original_height_emu) / 360000
        original_ratio = original_width_cm / original_height_cm
        logger.info(f"--- [PPTFiller]: 【尺寸检查】原始PPT尺寸:")
        logger.info(f"   宽度: {original_width_cm:.2f}cm ({original_width_emu} EMU)")
        logger.info(f"   高度: {original_height_cm:.2f}cm ({original_height_emu} EMU)")
        logger.info(f"   宽高比: {original_ratio:.2f}")
        logger.info(f"   16:9 = {16/9:.2f}, 4:3 = {4/3:.2f}")
        logger.info(f"   是否为16:9: {abs(original_ratio - 16/9) < 0.1}")
        logger.info(f"   是否为4:3: {abs(original_ratio - 4/3) < 0.1}")
        
        # 【修复1】强制设置为16:9
        target_width_cm = 33.867  # 16:9宽度
        target_height_cm = 19.05  # 16:9高度
        if abs(original_ratio - 16/9) > 0.1:
            logger.warning(f"--- [PPTFiller]: 【尺寸修复】检测到非16:9比例，正在转换为16:9...")
            prs.slide_width = Cm(target_width_cm)
            prs.slide_height = Cm(target_height_cm)
            new_ratio = target_width_cm / target_height_cm
            logger.info(f"--- [PPTFiller]: 【尺寸修复】已设置为16:9:")
            logger.info(f"   新宽度: {target_width_cm:.2f}cm")
            logger.info(f"   新高度: {target_height_cm:.2f}cm")
            logger.info(f"   新宽高比: {new_ratio:.2f} (目标: {16/9:.2f})")
        else:
            logger.info(f"--- [PPTFiller]: 【尺寸检查】PPT已经是16:9，无需修改")
        
        # 【日志探针2】检查设计规范应用
        logger.info(f"--- [PPTFiller]: 【设计规范】开始应用Ant Design设计规范...")
        logger.info(f"   主色: {ant_design_theme.colors.colorPrimary}")
        logger.info(f"   文本色: {ant_design_theme.colors.colorText}")
        logger.info(f"   字体族: {ant_design_theme.typography.fontFamily}")
        logger.info(f"   标题字号: {ant_design_theme.get_font_size_pt('h1')}pt")
        logger.info(f"   正文字号: {ant_design_theme.get_font_size_pt('base')}pt")
        
        # 构建样式策略（如果提供了分析结果）
        style_strategy = None
        if human_analysis and content_strategy:
            logger.info("--- [PPTFiller]: 【样式策略】构建智能样式策略...")
            style_strategy = self._build_style_strategy(human_analysis, content_strategy)
            logger.info(f"--- [PPTFiller]: 【样式策略】正式程度: {style_strategy.get('formality', '')}, 语调: {style_strategy.get('tone', '')}")
        
        # 填充每张幻灯片
        logger.info(f"--- [PPTFiller]: 【内容填充】开始填充 {len(prs.slides)} 张幻灯片...")
        for slide_idx, slide in enumerate(prs.slides):
            logger.info(f"--- [PPTFiller]: 【内容填充】处理幻灯片 {slide_idx + 1}/{len(prs.slides)}")
            placeholder_count = 0
            filled_count = 0
            
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholder_id = shape.placeholder_format.idx
                    key = f"slide_{slide_idx}_placeholder_{placeholder_id}"
                    placeholder_count += 1
                    
                    logger.debug(f"--- [PPTFiller]: 【占位符】幻灯片{slide_idx}, 占位符{placeholder_id}, key={key}")
                    
                    if key in content_map and hasattr(shape, "text_frame"):
                        try:
                            # 【日志探针3】记录填充前状态
                            old_text = shape.text_frame.text[:50] if shape.text_frame.text else "(空)"
                            logger.debug(f"--- [PPTFiller]: 【填充前】占位符{placeholder_id}内容: {old_text}...")
                            
                            # 清除现有内容
                            shape.text_frame.clear()
                            
                            # 添加新内容
                            content = content_map[key]
                            logger.debug(f"--- [PPTFiller]: 【内容】占位符{placeholder_id}新内容长度: {len(content)}字符")
                            
                            if content:
                                # 处理多段落
                                paragraphs = content.split('\n')
                                logger.debug(f"--- [PPTFiller]: 【段落】占位符{placeholder_id}包含{len(paragraphs)}个段落")
                                
                                for i, para_text in enumerate(paragraphs):
                                    if i == 0:
                                        p = shape.text_frame.paragraphs[0]
                                        p.text = para_text
                                    else:
                                        p = shape.text_frame.add_paragraph()
                                        p.text = para_text
                                
                                # 【改进】应用智能样式（如果提供了分析结果）
                                # 注意：对整个shape应用一次，而不是对每个段落
                                if style_strategy:
                                    self._apply_smart_style(
                                        shape, placeholder_id, slide_idx,
                                        style_strategy, human_analysis, content_strategy
                                    )
                                else:
                                    # 使用原有方法（向后兼容）
                                    for para in shape.text_frame.paragraphs:
                                        self._apply_ant_design_style(para, placeholder_id, slide_idx)
                                
                                filled_count += 1
                                logger.info(f"--- [PPTFiller]: 【填充成功】幻灯片{slide_idx}, 占位符{placeholder_id}")
                            else:
                                logger.warning(f"--- [PPTFiller]: 【内容为空】占位符{placeholder_id}没有内容")
                        except Exception as e:
                            logger.error(f"--- [PPTFiller]: 【填充失败】占位符{placeholder_id}: {e}", exc_info=True)
                    else:
                        if key not in content_map:
                            logger.warning(f"--- [PPTFiller]: 【缺少内容】占位符{placeholder_id}在content_map中不存在 (key={key})")
                        if not hasattr(shape, "text_frame"):
                            logger.warning(f"--- [PPTFiller]: 【无文本框架】占位符{placeholder_id}没有text_frame属性")
            
            logger.info(f"--- [PPTFiller]: 【幻灯片完成】幻灯片{slide_idx}: {filled_count}/{placeholder_count}个占位符已填充")
        
        logger.info(f"--- [PPTFiller]: 【内容填充】所有幻灯片处理完成")
        
        # 【日志探针4】最终尺寸检查
        final_width_emu = prs.slide_width
        final_height_emu = prs.slide_height
        final_width_cm = float(final_width_emu) / 360000
        final_height_cm = float(final_height_emu) / 360000
        final_ratio = final_width_cm / final_height_cm
        logger.info(f"--- [PPTFiller]: 【最终检查】保存前PPT尺寸:")
        logger.info(f"   宽度: {final_width_cm:.2f}cm ({final_width_emu} EMU)")
        logger.info(f"   高度: {final_height_cm:.2f}cm ({final_height_emu} EMU)")
        logger.info(f"   宽高比: {final_ratio:.2f} (目标16:9={16/9:.2f})")
        logger.info(f"   是否为16:9: {abs(final_ratio - 16/9) < 0.1}")
        
        # 保存
        prs.save(str(output_path))
        logger.info(f"--- [PPTFiller]: 【保存完成】PPT已保存到: {output_path}")
        
        # 【日志探针5】验证保存后的文件
        saved_size = Path(output_path).stat().st_size
        logger.info(f"--- [PPTFiller]: 【文件验证】保存后文件大小: {saved_size:,} bytes ({saved_size/1024:.2f} KB)")
        logger.info("="*80)
        logger.info("--- [PPTFiller]: 填充PPT流程完成")
        logger.info("="*80)
    
    def _build_style_strategy(
        self,
        human_analysis: Dict[str, Any],
        content_strategy: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        构建样式策略（充分利用所有6层分析结果）
        
        Args:
            human_analysis: 人类中心化分析结果
            content_strategy: 内容生成策略
            
        Returns:
            样式策略字典
        """
        # 【改进1】使用所有6层分析结果
        expression_style = human_analysis.get("layer_5_expression_style", {}).get("data", {})
        sections = human_analysis.get("layer_2_sections", {}).get("data", {})
        arguments = human_analysis.get("layer_3_arguments", {}).get("data", {})
        supporting_materials = human_analysis.get("layer_4_supporting_materials", {}).get("data", {})
        presentation_form = human_analysis.get("layer_6_presentation_form", {}).get("data", {})
        visual_style = content_strategy.get("expression_strategy", {}).get("visual_style", {})
        
        formality = expression_style.get("formality_level", "中性")
        tone = expression_style.get("tone", "中性")
        cultural_features = expression_style.get("cultural_features", [])
        
        # 【改进1.1】使用板块结构信息
        section_count = len(sections.get("sections", []))
        has_multiple_sections = section_count > 1
        
        # 【改进1.2】使用支撑材料信息
        has_data_points = len(supporting_materials.get("data_points", [])) > 0
        has_charts = len(supporting_materials.get("charts", [])) > 0
        has_case_studies = len(supporting_materials.get("case_studies", [])) > 0
        
        # 【改进1.3】使用呈现形式信息
        visual_hierarchy = presentation_form.get("visual_hierarchy", {})
        layout_style = presentation_form.get("layout_style", "standard")
        
        # 根据正式程度调整字号
        if formality == "正式":
            title_font_size = 40  # 更大，更正式
            body_font_size = 15
            subtitle_font_size = 24
        elif formality == "非正式":
            title_font_size = 36  # 稍小，更轻松
            body_font_size = 14
            subtitle_font_size = 22
        else:
            title_font_size = 38  # 标准
            body_font_size = 14
            subtitle_font_size = 24
        
        # 根据语调调整颜色
        if tone == "积极":
            primary_color = "#1890ff"  # Ant Design蓝色（积极）
            accent_color = "#52c41a"   # 绿色（成功）
            text_color = "#262626"     # 黑色
        elif tone == "谨慎":
            primary_color = "#fa8c16"  # 橙色（警告）
            accent_color = "#ff4d4f"    # 红色（错误）
            text_color = "#595959"     # 深灰色（更柔和）
        else:
            primary_color = "#1890ff"  # 标准蓝色
            accent_color = "#1890ff"
            text_color = "#262626"     # 标准黑色
        
        return {
            "typography": {
                "title_font_size": title_font_size,
                "subtitle_font_size": subtitle_font_size,
                "body_font_size": body_font_size,
                "font_family": visual_style.get("typography", {}).get("font_family", "Segoe UI")
            },
            "colors": {
                "primary": primary_color,
                "accent": accent_color,
                "text": text_color
            },
            "cultural_features": cultural_features,
            "formality": formality,
            "tone": tone,
            # 【改进1】新增：使用所有6层分析结果
            "sections": {
                "count": section_count,
                "has_multiple": has_multiple_sections
            },
            "supporting_materials": {
                "has_data_points": has_data_points,
                "has_charts": has_charts,
                "has_case_studies": has_case_studies
            },
            "presentation_form": {
                "visual_hierarchy": visual_hierarchy,
                "layout_style": layout_style
            }
        }
    
    def _determine_content_type(
        self,
        shape,
        placeholder_id: int,
        human_analysis: Optional[Dict[str, Any]]
    ) -> str:
        """
        确定内容类型
        
        Args:
            shape: PPT形状对象
            placeholder_id: 占位符ID
            human_analysis: 人类中心化分析结果
            
        Returns:
            内容类型字符串
        """
        import re
        
        # 获取占位符类型
        placeholder_type = ""
        if shape.is_placeholder:
            try:
                placeholder_type = str(shape.placeholder_format.type)
            except:
                pass
        
        # 获取文本内容
        text = ""
        if hasattr(shape, "text_frame"):
            text = shape.text_frame.text
        
        # 根据占位符类型判断
        if "CENTER_TITLE" in placeholder_type or "TITLE" in placeholder_type:
            return "title"
        elif "SUBTITLE" in placeholder_type:
            return "subtitle"
        elif "OBJECT" in placeholder_type or "BODY" in placeholder_type:
            # 检查内容是否包含数据或案例
            if "数据支撑" in text or re.search(r'\d+[%％]', text):
                return "data_highlight"
            elif "案例说明" in text or "案例" in text:
                return "case_study"
            elif "关键要点" in text or text.strip().startswith("•"):
                return "key_points"
            else:
                return "body"
        else:
            return "body"
    
    def _apply_smart_style(
        self,
        shape,
        placeholder_id: int,
        slide_idx: int,
        style_strategy: Dict[str, Any],
        human_analysis: Optional[Dict[str, Any]],
        content_strategy: Optional[Dict[str, Any]]
    ):
        """
        应用智能样式
        
        Args:
            shape: PPT形状对象
            placeholder_id: 占位符ID
            slide_idx: 幻灯片索引
            style_strategy: 样式策略
            human_analysis: 人类中心化分析结果
            content_strategy: 内容生成策略
        """
        # 1. 确定内容类型
        content_type = self._determine_content_type(shape, placeholder_id, human_analysis)
        logger.debug(f"--- [PPTFiller]: 【内容类型】占位符{placeholder_id}识别为: {content_type}")
        
        # 2. 根据内容类型应用样式
        if content_type == "title":
            self._apply_title_style(shape, style_strategy)
        elif content_type == "subtitle":
            self._apply_subtitle_style(shape, style_strategy)
        elif content_type == "data_highlight":
            self._apply_data_highlight_style(shape, style_strategy)
        elif content_type == "case_study":
            self._apply_case_study_style(shape, style_strategy)
        elif content_type == "key_points":
            self._apply_key_points_style(shape, style_strategy)
        else:
            self._apply_body_style(shape, style_strategy)
        
        # 3. 应用文化特征
        if "强调价值导向" in style_strategy.get("cultural_features", []):
            self._emphasize_value_propositions(shape, style_strategy)
        
        if "数据驱动表达" in style_strategy.get("cultural_features", []):
            self._emphasize_data_points(shape, style_strategy)
        
        # 【改进1】使用支撑材料信息突出数据
        supporting_materials = style_strategy.get("supporting_materials", {})
        if supporting_materials.get("has_data_points", False):
            self._emphasize_data_points(shape, style_strategy)
        
        # 【改进2】应用Ant Design间距系统
        self._apply_ant_design_spacing(shape, content_type)
        
        # 【改进2】应用Ant Design布局原则（增强视觉效果）
        self._apply_ant_design_layout(shape, content_type, style_strategy)
        
        # 【改进2.4】调整占位符位置（添加整体布局改进）
        self._adjust_placeholder_position(shape, content_type, slide_idx)
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """将hex颜色转换为RGBColor"""
        if hex_color.startswith('#'):
            r = int(hex_color[1:3], 16)
            g = int(hex_color[3:5], 16)
            b = int(hex_color[5:7], 16)
            return RGBColor(r, g, b)
        return RGBColor(38, 38, 38)  # 默认黑色
    
    def _apply_title_style(self, shape, style_strategy: Dict[str, Any]):
        """应用标题样式"""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                font.name = style_strategy["typography"]["font_family"]
                font.size = Pt(style_strategy["typography"]["title_font_size"])
                font.bold = True
                font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["primary"])
    
    def _apply_subtitle_style(self, shape, style_strategy: Dict[str, Any]):
        """应用副标题样式"""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                font.name = style_strategy["typography"]["font_family"]
                font.size = Pt(style_strategy["typography"]["subtitle_font_size"])
                font.bold = True
                font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["text"])
    
    def _apply_body_style(self, shape, style_strategy: Dict[str, Any]):
        """应用正文样式"""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                font.name = style_strategy["typography"]["font_family"]
                font.size = Pt(style_strategy["typography"]["body_font_size"])
                font.bold = False
                font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["text"])
    
    def _apply_data_highlight_style(self, shape, style_strategy: Dict[str, Any]):
        """应用数据高亮样式"""
        import re
        
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text
                font = run.font
                font.name = style_strategy["typography"]["font_family"]
                
                # 如果是数据（包含数字和%），使用强调色
                if re.search(r'\d+[%％]|\d+\.\d+%', text):
                    font.size = Pt(style_strategy["typography"]["body_font_size"] + 2)  # 稍大
                    font.bold = True
                    font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["accent"])
                else:
                    # 普通文本
                    font.size = Pt(style_strategy["typography"]["body_font_size"])
                    font.bold = False
                    font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["text"])
    
    def _apply_case_study_style(self, shape, style_strategy: Dict[str, Any]):
        """应用案例样式"""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                font.name = style_strategy["typography"]["font_family"]
                font.size = Pt(style_strategy["typography"]["body_font_size"])
                font.italic = True  # 案例使用斜体
                font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["text"])
    
    def _apply_key_points_style(self, shape, style_strategy: Dict[str, Any]):
        """应用关键要点样式"""
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                font.name = style_strategy["typography"]["font_family"]
                font.size = Pt(style_strategy["typography"]["body_font_size"])
                font.bold = True  # 要点加粗
                font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["text"])
    
    def _emphasize_value_propositions(self, shape, style_strategy: Dict[str, Any]):
        """突出价值主张"""
        import re
        
        value_keywords = ["降低", "提升", "加速", "优化", "改善", "增长", "提高"]
        
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if any(kw in run.text for kw in value_keywords):
                    # 价值主张使用强调色和加粗
                    run.font.bold = True
                    run.font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["accent"])
    
    def _emphasize_data_points(self, shape, style_strategy: Dict[str, Any]):
        """突出数据点"""
        import re
        
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text
                
                # 如果是数据，使用强调样式
                if re.search(r'\d+[%％]|\d+\.\d+%', text):
                    run.font.bold = True
                    run.font.size = Pt(style_strategy["typography"]["body_font_size"] + 2)
                    run.font.color.rgb = self._hex_to_rgb(style_strategy["colors"]["accent"])
    
    def _apply_ant_design_spacing(self, shape, content_type: str):
        """应用Ant Design间距系统"""
        from ant_design_theme import ant_design_theme
        from pptx.util import Pt
        
        try:
            # 获取Ant Design间距（转换为pt）
            spacing_sm_pt = ant_design_theme.get_spacing_cm('sm') * 28.35  # cm转pt (1cm ≈ 28.35pt)
            spacing_md_pt = ant_design_theme.get_spacing_cm('md') * 28.35
            spacing_lg_pt = ant_design_theme.get_spacing_cm('lg') * 28.35
            
            # 应用段落间距
            for i, para in enumerate(shape.text_frame.paragraphs):
                # 检查是否有paragraph_format属性
                if not hasattr(para, 'paragraph_format'):
                    continue
                
                try:
                    # 段落后间距（根据内容类型调整）
                    if content_type == "title":
                        para.paragraph_format.space_after = Pt(0)  # 标题后无间距
                    elif content_type == "subtitle":
                        para.paragraph_format.space_after = Pt(spacing_sm_pt)  # 副标题后小间距
                    elif content_type in ["data_highlight", "case_study", "key_points"]:
                        para.paragraph_format.space_after = Pt(spacing_md_pt)  # 重要内容后中等间距
                    else:
                        para.paragraph_format.space_after = Pt(spacing_sm_pt)  # 正文后小间距
                    
                    # 段落前间距（第一个段落除外）
                    if i > 0:
                        if content_type == "title":
                            para.paragraph_format.space_before = Pt(spacing_lg_pt)  # 标题前大间距
                        else:
                            para.paragraph_format.space_before = Pt(0)  # 其他内容前无间距
                except Exception as e:
                    logger.debug(f"--- [PPTFiller]: 【间距】应用段落间距失败: {e}")
            
            logger.debug(f"--- [PPTFiller]: 【间距】已应用Ant Design间距系统到{content_type}")
        except Exception as e:
            logger.warning(f"--- [PPTFiller]: 【间距】应用Ant Design间距系统失败: {e}")
    
    def _apply_ant_design_layout(self, shape, content_type: str, style_strategy: Dict[str, Any]):
        """应用Ant Design布局原则（增强视觉效果）"""
        from ant_design_theme import ant_design_theme
        from pptx.util import Cm
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        
        # 【改进2.1】应用文本对齐（根据内容类型）
        try:
            if content_type == "title":
                for para in shape.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER  # 标题居中
            elif content_type in ["data_highlight", "case_study", "key_points"]:
                for para in shape.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT  # 数据和案例左对齐
            else:
                for para in shape.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT  # 正文左对齐
        except Exception as e:
            logger.debug(f"--- [PPTFiller]: 【布局】应用对齐失败: {e}")
        
        # 【改进2.2】应用文本框架内边距（增强视觉效果）
        try:
            # 根据内容类型使用不同的内边距
            if content_type == "title":
                padding_cm = ant_design_theme.get_spacing_cm('lg')  # 24px = 更大间距
            elif content_type in ["data_highlight", "case_study", "key_points"]:
                padding_cm = ant_design_theme.get_spacing_cm('md')  # 16px = 中等间距
            else:
                padding_cm = ant_design_theme.get_spacing_cm('sm')  # 12px = 小间距
            
            # 应用内边距（通过调整文本框架的边距）
            text_frame = shape.text_frame
            text_frame.margin_left = Cm(padding_cm)
            text_frame.margin_right = Cm(padding_cm)
            text_frame.margin_top = Cm(padding_cm * 0.75)  # 上下间距稍大
            text_frame.margin_bottom = Cm(padding_cm * 0.75)
            
            logger.debug(f"--- [PPTFiller]: 【布局】已应用内边距: {padding_cm:.2f}cm ({content_type})")
        except Exception as e:
            logger.warning(f"--- [PPTFiller]: 【布局】应用内边距失败: {e}")
        
        # 【改进2.3】为所有内容类型添加背景色和边框（增强视觉效果）
        try:
            fill = shape.fill
            
            # 根据内容类型使用不同的背景色
            if content_type == "title":
                # 标题：使用主色背景（更明显）
                fill.solid()
                fill.fore_color.rgb = RGBColor(24, 144, 255)  # #1890ff 主色
                # 标题文字改为白色（在蓝色背景上）
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # 白色
            elif content_type in ["data_highlight", "case_study"]:
                # 数据和案例：使用浅灰背景
                fill.solid()
                fill.fore_color.rgb = RGBColor(240, 242, 245)  # #f0f2f5 浅灰
            elif content_type == "key_points":
                # 关键要点：使用极浅灰背景
                fill.solid()
                fill.fore_color.rgb = RGBColor(250, 250, 250)  # #fafafa 极浅灰
            else:
                # 正文：使用白色背景（保持原样）
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)  # #ffffff 白色
            
            # 添加边框（除了标题）
            if content_type != "title":
                try:
                    line = shape.line
                    line.color.rgb = RGBColor(217, 217, 217)  # #d9d9d9 边框色
                    line.width = Pt(1)  # 1pt边框
                except Exception as e:
                    logger.debug(f"--- [PPTFiller]: 【布局】应用边框失败: {e}")
            
            logger.debug(f"--- [PPTFiller]: 【布局】已应用背景色和边框到{content_type}")
        except Exception as e:
            logger.warning(f"--- [PPTFiller]: 【布局】应用背景色失败: {e}")
        
        logger.debug(f"--- [PPTFiller]: 【布局】已应用Ant Design布局原则到{content_type}")
    
    def _adjust_placeholder_position(self, shape, content_type: str, slide_idx: int):
        """调整占位符位置（添加整体布局改进）"""
        from ant_design_theme import ant_design_theme
        from pptx.util import Cm
        
        try:
            # 获取Ant Design间距
            slide_padding_cm = ant_design_theme.get_spacing_cm('lg')  # 24px = 0.63cm
            
            # 调整占位符位置（添加幻灯片内边距）
            # 注意：这会影响所有占位符，所以需要谨慎
            # 我们只调整位置，不改变尺寸（保持原有布局）
            
            # 如果占位符太靠近边缘，向内移动
            current_left_cm = float(shape.left) / 360000
            current_top_cm = float(shape.top) / 360000
            
            # 如果左边距小于内边距，调整位置
            if current_left_cm < slide_padding_cm:
                shape.left = Cm(slide_padding_cm)
                logger.debug(f"--- [PPTFiller]: 【布局】调整占位符左边距: {current_left_cm:.2f}cm -> {slide_padding_cm:.2f}cm")
            
            # 如果上边距小于内边距，调整位置
            if current_top_cm < slide_padding_cm:
                shape.top = Cm(slide_padding_cm)
                logger.debug(f"--- [PPTFiller]: 【布局】调整占位符上边距: {current_top_cm:.2f}cm -> {slide_padding_cm:.2f}cm")
            
        except Exception as e:
            logger.debug(f"--- [PPTFiller]: 【布局】调整占位符位置失败: {e}")
    
    def _apply_ant_design_style(self, paragraph, placeholder_id: int, slide_idx: int):
        """
        应用Ant Design设计规范到段落
        
        Args:
            paragraph: PPT段落对象
            placeholder_id: 占位符ID
            slide_idx: 幻灯片索引
        """
        try:
            if not paragraph.runs:
                run = paragraph.add_run()
            else:
                run = paragraph.runs[0]
            
            font = run.font
            
            # 【修复3】应用Ant Design字体系统
            logger.debug(f"--- [PPTFiller]: 【字体应用】幻灯片{slide_idx}, 占位符{placeholder_id}")
            try:
                font.name = "Segoe UI"  # Windows系统字体
                logger.debug(f"--- [PPTFiller]: 【字体】设置为: Segoe UI")
            except Exception:
                try:
                    font.name = "Helvetica Neue"  # macOS系统字体
                    logger.debug(f"--- [PPTFiller]: 【字体】设置为: Helvetica Neue")
                except Exception:
                    try:
                        font.name = "微软雅黑"  # 中文字体fallback
                        logger.debug(f"--- [PPTFiller]: 【字体】设置为: 微软雅黑")
                    except Exception:
                        font.name = "Arial"  # 最终fallback
                        logger.debug(f"--- [PPTFiller]: 【字体】设置为: Arial (fallback)")
            
            # 【修复4】应用Ant Design字号系统
            # 根据占位符类型判断（通常placeholder_id=0是标题，1是内容）
            if placeholder_id == 0:
                # 标题
                font.size = Pt(ant_design_theme.get_font_size_pt('h1'))
                font.bold = True
                logger.debug(f"--- [PPTFiller]: 【字号】占位符{placeholder_id}设置为标题: {ant_design_theme.get_font_size_pt('h1')}pt, 加粗")
            else:
                # 正文
                font.size = Pt(ant_design_theme.get_font_size_pt('base'))
                font.bold = False
                logger.debug(f"--- [PPTFiller]: 【字号】占位符{placeholder_id}设置为正文: {ant_design_theme.get_font_size_pt('base')}pt")
            
            # 【修复5】应用Ant Design颜色系统
            try:
                # 解析颜色（从hex转换为RGB）
                text_color_hex = ant_design_theme.colors.colorText
                if text_color_hex.startswith('#'):
                    r = int(text_color_hex[1:3], 16)
                    g = int(text_color_hex[3:5], 16)
                    b = int(text_color_hex[5:7], 16)
                    font.color.rgb = RGBColor(r, g, b)
                    logger.debug(f"--- [PPTFiller]: 【颜色】设置为: {text_color_hex} (RGB: {r}, {g}, {b})")
            except Exception as e:
                logger.warning(f"--- [PPTFiller]: 【颜色应用失败】占位符{placeholder_id}: {e}")
            
        except Exception as e:
            logger.error(f"--- [PPTFiller]: 【样式应用失败】占位符{placeholder_id}: {e}", exc_info=True)
    
    async def _generate_content_by_sections(
        self,
        human_analysis: Dict[str, Any],
        content_strategy: Dict[str, Any],
        user_prompt: str,
        intelligent_data_points: Optional[List[Dict[str, Any]]] = None,
        intelligent_cases: Optional[List[Dict[str, Any]]] = None
    ) -> Dict[str, str]:
        """
        逐板块生成内容（改进方案）
        
        Args:
            human_analysis: 人类中心化分析结果
            content_strategy: 内容生成策略
            user_prompt: 用户需求
            
        Returns:
            内容映射字典
        """
        content_map = {}
        section_strategies = content_strategy.get("section_strategies", [])
        
        logger.info(f"--- [PPTFiller]: Generating content for {len(section_strategies)} sections...")
        
        # 初始化润色器、展示策划器、布局规划器和颜色配置器
        content_polisher = ContentPolisher(self.llm_service)
        presentation_planner = PresentationPlanner(self.llm_service)
        layout_planner = LayoutPlanner(self.llm_service)
        from ..theme.color_configurator import ColorConfigurator
        color_configurator = ColorConfigurator(self.llm_service)
        
        # 【新增】收集所有板块的润色结果、布局规划结果和颜色配置结果
        all_polished_slides = []
        all_presentation_plans = []
        all_layout_plans = []
        all_color_configs = []
        
        for idx, section_strategy in enumerate(section_strategies):
            logger.info(f"--- [PPTFiller]: Processing section {idx + 1}/{len(section_strategies)}: {section_strategy.get('theme', '')}")
            
            # 获取该板块相关的支撑材料
            section_slides = section_strategy.get("slides", [])
            section_data_points = []
            section_cases = []
            
            if intelligent_data_points:
                section_data_points = [
                    dp for dp in intelligent_data_points 
                    if dp.get("slide_index", -1) in section_slides
                ]
            
            if intelligent_cases:
                section_cases = [
                    c for c in intelligent_cases 
                    if c.get("slide_index", -1) in section_slides
                ]
            
            logger.info(f"--- [PPTFiller]: Section {idx + 1} 相关支撑材料: {len(section_data_points)}个数据点, {len(section_cases)}个案例")
            
            # 【新增】步骤0: 获取板块分析结果（用于润色）
            section_analysis = self._get_section_analysis(idx, human_analysis, section_strategy)
            
            # 【新增】步骤1: 内容润色 - 将文档内容润色成适合PPT展示的文案
            logger.info(f"--- [PPTFiller]: 【步骤1】内容润色 - 板块{idx + 1}...")
            polished_slides = await content_polisher.polish_section(
                section_analysis=section_analysis,
                section_index=idx
            )
            logger.info(f"--- [PPTFiller]: ✅ 润色完成，生成{len(polished_slides)}张幻灯片")
            
            # 【新增】步骤2: 展示策划 - 设计简洁明了的展示方式
            logger.info(f"--- [PPTFiller]: 【步骤2】展示策划 - 板块{idx + 1}...")
            presentation_plan = await presentation_planner.plan_presentation(
                polished_slides=polished_slides,
                section_theme=section_strategy.get('theme', '')
            )
            logger.info(f"--- [PPTFiller]: ✅ 展示策划完成，策划了{len(presentation_plan)}张幻灯片")
            
            # 【新增】步骤3: 布局规划 - 基于设计规范生成详细的布局说明
            logger.info(f"--- [PPTFiller]: 【步骤3】布局规划 - 板块{idx + 1}...")
            layout_plans = await layout_planner.plan_layout(
                polished_slides=polished_slides,
                presentation_plan=presentation_plan
            )
            logger.info(f"--- [PPTFiller]: ✅ 布局规划完成，规划了{len(layout_plans)}张幻灯片")
            
            # 【新增】步骤4: 颜色配置 - 为每个元素配置符合Ant Design规范的颜色方案
            logger.info(f"--- [PPTFiller]: 【步骤4】颜色配置 - 板块{idx + 1}...")
            color_configs = await color_configurator.configure_colors(
                polished_slides=polished_slides,
                presentation_plans=presentation_plan,
                layout_plans=layout_plans
            )
            logger.info(f"--- [PPTFiller]: ✅ 颜色配置完成，配置了{len(color_configs)}张幻灯片")
            
            # 【新增】收集润色结果、展示策划结果、布局规划结果和颜色配置结果
            # 【探针】记录板块内的slide_index分布
            logger.info(f"--- [PPTFiller]: 【探针】板块{idx + 1}的slide_index分布:")
            logger.info(f"--- [PPTFiller]:   polished_slides: {[s.get('slide_index', 0) for s in polished_slides]}")
            logger.info(f"--- [PPTFiller]:   layout_plans: {[l.get('slide_index', 0) for l in layout_plans]}")
            
            # 在合并前，调整slide_index为全局索引
            # 计算全局起始索引
            global_start_index = sum(len(all_polished_slides) for _ in range(idx)) if idx > 0 else 0
            global_start_index = len(all_polished_slides)  # 使用实际已收集的数量
            
            logger.info(f"--- [PPTFiller]: 【探针】板块{idx + 1}的全局起始索引: {global_start_index}")
            
            # 调整polished_slides的slide_index
            adjusted_polished_slides = []
            for slide in polished_slides:
                adjusted_slide = slide.copy()
                old_index = adjusted_slide.get('slide_index', 0)
                new_index = global_start_index + old_index
                adjusted_slide['slide_index'] = new_index
                adjusted_slide['section_index'] = idx  # 保留板块索引
                adjusted_polished_slides.append(adjusted_slide)
                logger.debug(f"--- [PPTFiller]:   调整polished_slide: {old_index} -> {new_index}")
            
            # 调整layout_plans的slide_index
            adjusted_layout_plans = []
            for plan in layout_plans:
                adjusted_plan = plan.copy()
                old_index = adjusted_plan.get('slide_index', 0)
                new_index = global_start_index + old_index
                adjusted_plan['slide_index'] = new_index
                adjusted_plan['section_index'] = idx  # 保留板块索引
                adjusted_layout_plans.append(adjusted_plan)
                logger.debug(f"--- [PPTFiller]:   调整layout_plan: {old_index} -> {new_index}")
            
            # 调整presentation_plans的slide_index
            adjusted_presentation_plans = []
            for plan in presentation_plan:
                adjusted_plan = plan.copy()
                old_index = adjusted_plan.get('slide_index', 0)
                new_index = global_start_index + old_index
                adjusted_plan['slide_index'] = new_index
                adjusted_plan['section_index'] = idx  # 保留板块索引
                adjusted_presentation_plans.append(adjusted_plan)
                logger.debug(f"--- [PPTFiller]:   调整presentation_plan: {old_index} -> {new_index}")
            
            # 调整color_configs的slide_index
            adjusted_color_configs = []
            for config in color_configs:
                adjusted_config = config.copy()
                old_index = adjusted_config.get('slide_index', 0)
                new_index = global_start_index + old_index
                adjusted_config['slide_index'] = new_index
                adjusted_config['section_index'] = idx  # 保留板块索引
                adjusted_color_configs.append(adjusted_config)
                logger.debug(f"--- [PPTFiller]:   调整color_config: {old_index} -> {new_index}")
            
            # 使用调整后的数据
            all_polished_slides.extend(adjusted_polished_slides)
            all_presentation_plans.extend(adjusted_presentation_plans)
            all_layout_plans.extend(adjusted_layout_plans)
            all_color_configs.extend(adjusted_color_configs)
            
            logger.info(f"--- [PPTFiller]: 【探针】合并后总数: polished_slides={len(all_polished_slides)}, layout_plans={len(all_layout_plans)}")
            
            # 1. 构建板块特定的生成提示词（包含支撑材料、润色结果、展示策划）
            section_prompt = self._build_section_prompt(
                section_strategy=section_strategy,
                overall_strategy=content_strategy.get("overall_strategy", {}),
                expression_strategy=content_strategy.get("expression_strategy", {}),
                user_prompt=user_prompt,
                human_analysis=human_analysis,
                section_strategies=section_strategies,
                data_points=section_data_points,
                cases=section_cases,
                polished_slides=polished_slides,
                presentation_plan=presentation_plan
            )
            
            # 2. 调用LLM生成板块内容（基于润色和展示策划结果）
            logger.info(f"--- [PPTFiller]: 【详细探针】调用LLM生成板块{idx + 1}内容...")
            logger.info(f"   提示词长度: {len(section_prompt)}字符")
            logger.info(f"   提示词预览: {section_prompt[:500]}...")
            
            section_content = await self._generate_section_content(section_prompt)
            
            logger.info(f"--- [PPTFiller]: 【详细探针】LLM生成结果:")
            logger.info(f"   返回类型: {type(section_content)}")
            if isinstance(section_content, dict):
                logger.info(f"   标题: {section_content.get('title', '')[:100]}...")
                logger.info(f"   内容长度: {len(section_content.get('content', ''))}字符")
                logger.info(f"   关键点数量: {len(section_content.get('key_points', []))}")
                logger.info(f"   数据高亮数量: {len(section_content.get('data_highlights', []))}")
                logger.info(f"   案例数量: {len(section_content.get('case_studies', []))}")
            else:
                logger.warning(f"   ⚠️ 返回格式异常: {section_content}")
            
            # 3. 结构化板块内容
            logger.info(f"--- [PPTFiller]: 【详细探针】结构化板块内容...")
            structured_content = self._structure_section_content(
                section_content,
                section_strategy
            )
            
            logger.info(f"--- [PPTFiller]: 【详细探针】结构化结果:")
            logger.info(f"   标题: {structured_content.get('title', '')[:100]}...")
            logger.info(f"   主要内容: {len(structured_content.get('main_content', ''))}字符")
            logger.info(f"   关键点: {len(structured_content.get('key_points', []))}个")
            if structured_content.get('key_points'):
                for i, kp in enumerate(structured_content['key_points'][:3]):
                    logger.info(f"     关键点{i+1}: {kp[:80]}...")
            
            # 4. 映射到PPT占位符
            section_slides = section_strategy.get("slides", [])
            logger.info(f"--- [PPTFiller]: 【探针】板块{idx + 1}映射到幻灯片: {section_slides}")
            section_content_map = self._map_to_placeholders(
                structured_content,
                section_slides,
                human_analysis
            )
            
            # 【探针】记录映射结果
            logger.info(f"--- [PPTFiller]: 【探针】板块{idx + 1}生成的内容映射:")
            for key, content in section_content_map.items():
                logger.info(f"   {key}: {len(content)}字符 - {content[:50]}...")
            
            # 5. 合并到总内容映射
            content_map.update(section_content_map)
            logger.info(f"--- [PPTFiller]: Section {idx + 1} completed, generated {len(section_content_map)} placeholders")
            logger.info(f"--- [PPTFiller]: 【探针】当前总内容映射项数: {len(content_map)}")
        
        # 【新增】返回内容映射、润色结果、布局规划结果和颜色配置结果
        return {
            'content_map': content_map,
            'polished_slides': all_polished_slides,
            'presentation_plans': all_presentation_plans,
            'layout_plans': all_layout_plans,
            'color_configs': all_color_configs
        }
    
    def _get_section_analysis(
        self,
        section_index: int,
        human_analysis: Dict[str, Any],
        section_strategy: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        获取特定板块的分析结果（用于润色）
        
        Args:
            section_index: 板块索引
            human_analysis: 人类中心化分析结果
            section_strategy: 板块策略
            
        Returns:
            板块分析结果字典
        """
        sections = human_analysis.get("layer_2_sections", {}).get("data", {})
        section_list = sections.get("sections", [])
        
        # 查找对应板块
        section_data = None
        for section in section_list:
            if section.get("section_index", -1) == section_index:
                section_data = section
                break
        
        # 如果没有找到，使用板块策略中的信息
        if not section_data:
            section_data = {
                "section_index": section_index,
                "theme": section_strategy.get("theme", ""),
                "core_idea": section_strategy.get("core_idea", ""),
                "content_summary": section_strategy.get("content_summary", "")
            }
        
        # 获取该板块的论证逻辑
        arguments = human_analysis.get("layer_3_arguments", {}).get("data", {})
        argument_list = arguments.get("arguments", [])
        section_arguments = [
            arg for arg in argument_list 
            if arg.get("section_index", -1) == section_index
        ]
        
        return {
            "section_index": section_index,
            "theme": section_data.get("theme", ""),
            "core_idea": section_data.get("core_idea", ""),
            "content_summary": section_data.get("content_summary", ""),
            "arguments": section_arguments,
            "slides": section_strategy.get("slides", [])
        }
    
    def _build_section_prompt(
        self,
        section_strategy: Dict[str, Any],
        overall_strategy: Dict[str, Any],
        expression_strategy: Dict[str, Any],
        user_prompt: str,
        human_analysis: Dict[str, Any],
        section_strategies: List[Dict[str, Any]],
        data_points: Optional[List[Dict[str, Any]]] = None,
        cases: Optional[List[Dict[str, Any]]] = None,
        polished_slides: Optional[List[Dict[str, Any]]] = None,
        presentation_plan: Optional[List[Dict[str, Any]]] = None
    ) -> str:
        """构建板块特定的生成提示词"""
        
        section_idx = section_strategy.get("section_index", 0)
        approach = section_strategy.get("content_generation_approach", {})
        
        # 获取前后板块信息
        prev_section = section_strategies[section_idx - 1] if section_idx > 0 else None
        next_section = section_strategies[section_idx + 1] if section_idx < len(section_strategies) - 1 else None
        
        # 格式化证据点
        evidence_points = section_strategy.get("evidence_points", [])
        evidence_text = "\n".join([f"{i+1}. {point}" for i, point in enumerate(evidence_points[:5])]) if evidence_points else "无"
        
        prompt = f"""
【整体背景】
核心主题：{overall_strategy.get("core_theme", "")}
价值主张：{", ".join(overall_strategy.get("value_propositions", []))}
目标受众：{overall_strategy.get("target_audience", "")}
文档目的：{overall_strategy.get("purpose", "")}

【当前板块】
板块主题：{section_strategy.get("theme", "")}
核心思想：{section_strategy.get("core_idea", "")}
板块位置：第{section_idx + 1}个板块（共{len(section_strategies)}个）

【论证逻辑】
论证方式：{", ".join(section_strategy.get("argument_types", [])) if section_strategy.get("argument_types") else "通用论证"}
证据点：
{evidence_text}

【生成策略】
强调重点：{approach.get("emphasis", "核心观点")}
证据优先级：{", ".join(approach.get("evidence_priority", []))}
内容长度：{approach.get("length", "medium")}

【表达风格】
正式程度：{expression_strategy.get("language_style", {}).get("formality", "中性")}
语调：{expression_strategy.get("language_style", {}).get("tone", "中性")}
文化特征：{", ".join(expression_strategy.get("language_style", {}).get("cultural_features", []))}

【板块上下文】
前一个板块：{prev_section.get("theme", "无") if prev_section else "无"}
后一个板块：{next_section.get("theme", "无") if next_section else "无"}

【可用支撑材料】
{self._format_supporting_materials_for_prompt(data_points, cases)}

【润色结果】
{self._format_polished_slides_for_prompt(polished_slides) if polished_slides else "未进行润色"}

【展示策划】
{self._format_presentation_plan_for_prompt(presentation_plan) if presentation_plan else "未进行展示策划"}

【用户需求】
{user_prompt}

【生成要求 - 重要：这是PPT，不是发言稿！】
1. **标题**：简洁有力，体现核心思想，不超过15字
2. **正文（content字段）**：只放简洁概述（1-2句话，不超过50字），不要放长段落
3. **关键要点（key_points）**：必须是要点式列表，每个要点不超过30字，不要用完整句子
4. **数据**：优先使用上述"可用数据点"，如果有数据论证，突出显示数据（如"降低40-60%成本"）
5. **案例**：优先使用上述"可用案例"，如果有案例论证，简洁说明案例（要点式，不超过20字）
6. **逻辑**：与前一个板块有逻辑衔接，为后一个板块做铺垫
7. **风格**：符合{expression_strategy.get("language_style", {}).get("formality", "中性")}风格，语调{expression_strategy.get("language_style", {}).get("tone", "中性")}
8. **文化**：体现{", ".join(expression_strategy.get("language_style", {}).get("cultural_features", [])) if expression_strategy.get("language_style", {}).get("cultural_features") else "通用商业文化"}

【禁止事项】：
- ❌ 不要生成"各位同事"、"今天我将"等发言稿开场白
- ❌ 不要生成长段落描述，必须用要点列表
- ❌ 不要生成完整的句子，要用简洁的要点
- ❌ content字段不要放长段落，应该放简洁的概述（1-2句话），详细内容放在key_points中

【正确格式示例】：
{{
  "title": "技术产品商业化全链路AI解决方案",
  "content": "基于25年技术积累，形成完整的AI赋能解决方案体系",
  "key_points": [
    "降低运营成本40-60%",
    "提升转化效率20-35%",
    "三大核心系统：朋友云、BefriendsAI、数据中心"
  ],
  "data_highlights": ["降低运营成本40-60%", "提升转化效率20-35%"],
  "case_studies": []
}}

请生成该板块的内容（JSON格式）：
{{
  "title": "标题（不超过15字）",
  "content": "简洁概述（1-2句话，不超过50字）",
  "key_points": ["要点1（不超过30字）", "要点2（不超过30字）"],
  "data_highlights": ["数据1", "数据2"],
  "case_studies": ["案例1（不超过20字）", "案例2（不超过20字）"]
}}
"""
        return prompt
    
    async def _generate_section_content(self, section_prompt: str) -> Dict[str, Any]:
        """生成板块内容"""
        system_prompt = """你是一个专业的PPT内容创作助手，擅长生成符合中国商业汇报习惯的PPT内容。

【重要】PPT内容要求（不是发言稿！）：
1. **简洁性**：PPT是视觉辅助工具，内容必须简洁明了，每页不超过5-7个要点
2. **要点式**：使用要点列表（bullet points），而不是长段落描述
3. **标题**：标题要简洁有力，不超过15字，体现核心观点
4. **正文**：正文应该是要点列表，每个要点不超过30字，避免完整句子
5. **数据**：如果有数据，单独列出，突出显示（如"降低40-60%成本"）
6. **避免**：
   - ❌ 不要生成完整的发言稿或演讲稿
   - ❌ 不要使用"各位同事"、"今天我将"等开场白
   - ❌ 不要生成长段落描述
   - ❌ 不要使用完整的句子，要用要点式
7. **正确示例**：
   - ✅ "降低运营成本40-60%"
   - ✅ "提升转化效率20-35%"
   - ✅ "三大核心系统：朋友云、BefriendsAI、数据中心"
   - ❌ "各位管理层同事，今天我将向大家汇报..."（这是发言稿，不是PPT）

要求：
1. 理解板块的核心思想和论证逻辑
2. 根据生成策略组织内容，但必须是要点式，不是详细描述
3. 突出价值主张和关键数据
4. 保持逻辑清晰，有说服力
5. 符合表达风格和文化特征

输出必须是有效的JSON格式。"""
        
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": section_prompt}
        ]
        
        try:
            logger.debug(f"--- [PPTFiller]: 【详细探针】调用LLM API...")
            logger.debug(f"   模型: {self.llm_service.model_name}")
            logger.debug(f"   Base URL: {self.llm_service.base_url}")
            logger.debug(f"   消息数量: {len(messages)}")
            logger.debug(f"   System消息长度: {len(messages[0].get('content', ''))}字符")
            logger.debug(f"   User消息长度: {len(messages[1].get('content', ''))}字符")
            
            response = await self.llm_service.chat_completion_async(
                messages=messages,
                temperature=0.7,
                max_tokens=2000
            )
            
            logger.debug(f"--- [PPTFiller]: 【详细探针】LLM API调用成功")
            logger.debug(f"   响应长度: {len(response)}字符")
            logger.debug(f"   响应预览: {response[:300]}...")
            
            # 解析JSON响应
            import json
            import re
            
            # 尝试提取JSON（支持多种格式）
            json_match = re.search(r'\{.*\}', response, re.DOTALL)
            if json_match:
                try:
                    parsed_json = json.loads(json_match.group(0))
                    logger.debug(f"--- [PPTFiller]: 【详细探针】JSON解析成功")
                    logger.debug(f"   解析后的键: {list(parsed_json.keys())}")
                    return parsed_json
                except json.JSONDecodeError as je:
                    logger.error(f"--- [PPTFiller]: JSON解析失败: {je}")
                    logger.error(f"   尝试解析的内容: {json_match.group(0)[:500]}...")
                    # 尝试修复常见的JSON问题
                    json_str = json_match.group(0)
                    # 移除可能的注释
                    json_str = re.sub(r'//.*?$', '', json_str, flags=re.MULTILINE)
                    try:
                        return json.loads(json_str)
                    except:
                        pass
            else:
                logger.warning("--- [PPTFiller]: 未找到JSON格式，使用fallback")
                logger.warning(f"   响应内容: {response[:500]}...")
            
            # Fallback: 尝试从响应中提取内容
            return {
                "title": "标题",
                "content": response[:200] if response else "",
                "key_points": [],
                "data_highlights": [],
                "case_studies": []
            }
        except Exception as e:
            logger.error(f"--- [PPTFiller]: Failed to generate section content: {e}", exc_info=True)
            logger.error(f"--- [PPTFiller]: 【详细探针】错误类型: {type(e).__name__}")
            logger.error(f"--- [PPTFiller]: 【详细探针】错误详情: {str(e)}")
            
            # 如果是API错误，记录更多信息
            if hasattr(e, 'response'):
                logger.error(f"--- [PPTFiller]: 【详细探针】API响应: {e.response}")
            if hasattr(e, 'status_code'):
                logger.error(f"--- [PPTFiller]: 【详细探针】HTTP状态码: {e.status_code}")
            
            return {
                "title": "标题",
                "content": "",
                "key_points": [],
                "data_highlights": [],
                "case_studies": []
            }
    
    def _structure_section_content(
        self,
        section_content: Dict[str, Any],
        section_strategy: Dict[str, Any]
    ) -> Dict[str, Any]:
        """结构化板块内容"""
        structured = {
            "title": section_content.get("title", ""),
            "main_content": section_content.get("content", ""),
            "key_points": section_content.get("key_points", []),
            "data_highlights": section_content.get("data_highlights", []),
            "case_studies": section_content.get("case_studies", []),
            "section_index": section_strategy.get("section_index", 0),
            "slides": section_strategy.get("slides", [])
        }
        
        # 根据论证类型调整结构
        argument_types = section_strategy.get("argument_types", [])
        if "数据论证" in argument_types:
            structured["emphasis"] = "data"
            structured["data_formatted"] = self._format_data_highlights(
                structured["data_highlights"]
            )
        
        if "案例论证" in argument_types:
            structured["emphasis"] = "case"
            structured["cases_formatted"] = self._format_case_studies(
                structured["case_studies"]
            )
        
        return structured
    
    def _format_polished_slides_for_prompt(
        self,
        polished_slides: List[Dict[str, Any]]
    ) -> str:
        """
        格式化润色后的幻灯片用于提示词
        
        Args:
            polished_slides: 润色后的幻灯片列表
            
        Returns:
            格式化后的字符串
        """
        if not polished_slides:
            return "未进行润色"
        
        formatted = []
        for slide in polished_slides:
            slide_idx = slide.get('slide_index', 0)
            title = slide.get('title', '')
            content = slide.get('content', '')
            content_type = slide.get('content_type', '')
            
            formatted.append(f"幻灯片{slide_idx} ({content_type}):")
            if title:
                formatted.append(f"  标题: {title}")
            if content:
                # 限制长度
                content_preview = content[:200] + "..." if len(content) > 200 else content
                formatted.append(f"  内容: {content_preview}")
            
            # 如果有视觉元素详情
            visual_elements = slide.get('visual_elements_detail', [])
            if visual_elements:
                formatted.append(f"  视觉元素 ({len(visual_elements)}个):")
                for elem in visual_elements[:3]:  # 只显示前3个
                    elem_id = elem.get('element_id', '')
                    elem_type = elem.get('element_type', '')
                    elem_title = elem.get('title', '')
                    formatted.append(f"    - {elem_id} ({elem_type}): {elem_title}")
        
        return "\n".join(formatted)
    
    def _format_presentation_plan_for_prompt(
        self,
        presentation_plan: List[Dict[str, Any]]
    ) -> str:
        """
        格式化展示策划结果用于提示词
        
        Args:
            presentation_plan: 展示策划结果列表
            
        Returns:
            格式化后的字符串
        """
        if not presentation_plan:
            return "未进行展示策划"
        
        formatted = []
        for plan in presentation_plan:
            slide_idx = plan.get('slide_index', 0)
            layout_type = plan.get('layout_type', '')
            layout_desc = plan.get('layout_description', '')
            visual_guidance = plan.get('visual_guidance', {})
            
            formatted.append(f"幻灯片{slide_idx}:")
            formatted.append(f"  布局类型: {layout_type}")
            if layout_desc:
                formatted.append(f"  布局描述: {layout_desc[:150]}...")
            if visual_guidance:
                font_size = visual_guidance.get('font_size', '')
                alignment = visual_guidance.get('alignment', '')
                if font_size or alignment:
                    formatted.append(f"  视觉指导: 字号{font_size}, 对齐{alignment}")
        
        return "\n".join(formatted)
    
    def _format_data_highlights(self, data_highlights: List[str]) -> str:
        """格式化数据高亮"""
        if not data_highlights:
            return ""
        
        formatted = []
        for data in data_highlights:
            formatted.append(f"• {data}")
        
        return "\n".join(formatted)
    
    def _format_case_studies(self, case_studies: List[str]) -> str:
        """格式化案例说明"""
        if not case_studies:
            return ""
        
        formatted = []
        for case in case_studies:
            formatted.append(f"• {case}")
        
        return "\n".join(formatted)
    
    def _format_supporting_materials_for_prompt(
        self,
        data_points: Optional[List[Dict[str, Any]]],
        cases: Optional[List[Dict[str, Any]]]
    ) -> str:
        """
        格式化支撑材料用于提示词
        
        Args:
            data_points: 智能识别的数据点列表
            cases: 智能识别的案例列表
            
        Returns:
            格式化的支撑材料文本
        """
        parts = []
        
        # 格式化数据点
        if data_points:
            parts.append("【可用数据点】")
            for dp in data_points:
                data_text = f"- {dp.get('value', '')}"
                if dp.get('label'):
                    data_text += f" ({dp.get('label')})"
                if dp.get('context'):
                    data_text += f" - {dp.get('context')[:50]}"
                parts.append(data_text)
        else:
            parts.append("【可用数据点】无")
        
        # 格式化案例
        if cases:
            parts.append("\n【可用案例】")
            for case in cases:
                case_text = f"- {case.get('type', '案例')}"
                if case.get('company'):
                    case_text += f": {case.get('company')}"
                if case.get('result'):
                    case_text += f" - 结果: {case.get('result')}"
                if case.get('content'):
                    case_text += f" - {case.get('content')[:100]}"
                parts.append(case_text)
        else:
            parts.append("\n【可用案例】无")
        
        return "\n".join(parts)
    
    def _map_to_placeholders(
        self,
        structured_content: Dict[str, Any],
        slide_indices: List[int],
        human_analysis: Dict[str, Any]
    ) -> Dict[str, str]:
        """智能映射到PPT占位符（详细探针）"""
        content_map = {}
        
        logger.info("="*80)
        logger.info("--- [PPTFiller]: 【详细探针】内容映射到PPT占位符")
        logger.info("="*80)
        logger.info(f"   结构化内容:")
        logger.info(f"     标题: {structured_content.get('title', '')[:100]}...")
        logger.info(f"     主要内容长度: {len(structured_content.get('main_content', ''))}字符")
        logger.info(f"     关键点数量: {len(structured_content.get('key_points', []))}")
        logger.info(f"     数据高亮数量: {len(structured_content.get('data_highlights', []))}")
        logger.info(f"     案例数量: {len(structured_content.get('case_studies', []))}")
        logger.info(f"   目标幻灯片: {slide_indices}")
        
        # 获取所有幻灯片的占位符信息
        slides_data = human_analysis.get("layer_1_physical", {}).get("data", {})
        # 从原始结构获取占位符信息
        enhanced_parser = EnhancedPPTParser(str(self.framework_path))
        enhanced_structure = enhanced_parser.extract_structure_enhanced()
        
        logger.info(f"   框架PPT幻灯片数: {len(enhanced_structure.get('slides', []))}")
        
        for slide_idx in slide_indices:
            logger.info(f"\n   --- 处理幻灯片{slide_idx} ---")
            if slide_idx < len(enhanced_structure.get("slides", [])):
                # 使用框架PPT的占位符结构
                slide_data = enhanced_structure["slides"][slide_idx]
                placeholders = slide_data.get("placeholders", [])
                
                logger.info(f"     占位符数量: {len(placeholders)}")
                
                # 根据占位符类型分配内容
                for placeholder in placeholders:
                    placeholder_id = placeholder.get("placeholder_id", 0)
                    placeholder_type = placeholder.get("placeholder_type", "")
                    placeholder_key = f"slide_{slide_idx}_placeholder_{placeholder_id}"
                    
                    logger.info(f"     占位符{placeholder_id}:")
                    logger.info(f"       类型: {placeholder_type}")
                    logger.info(f"       键名: {placeholder_key}")
                    
                    if "TITLE" in placeholder_type or "CENTER_TITLE" in placeholder_type:
                        # 标题占位符
                        assigned_content = structured_content["title"]
                        content_map[placeholder_key] = assigned_content
                        logger.info(f"       → 分配标题内容: {assigned_content[:100]}...")
                    elif "OBJECT" in placeholder_type or "BODY" in placeholder_type or "SUBTITLE" in placeholder_type:
                        # 正文占位符
                        # 【改进】根据占位符ID和内容类型智能分配
                        # 如果有多个占位符，可以拆分不同类型的内容
                        placeholder_idx = placeholder_id
                        
                        if placeholder_idx == 1:
                            # 第一个正文占位符：优先显示关键要点
                            if structured_content.get("key_points"):
                                key_points_text = "\n".join([f"• {point}" for point in structured_content["key_points"]])
                                content_map[placeholder_key] = key_points_text
                                logger.info(f"       → 分配关键要点: {len(key_points_text)}字符, {len(structured_content['key_points'])}个要点")
                            else:
                                content = self._build_body_content(structured_content, placeholder_type)
                                content_map[placeholder_key] = content
                                logger.info(f"       → 分配正文内容: {len(content)}字符")
                        elif placeholder_idx == 2:
                            # 第二个正文占位符：显示数据高亮
                            if structured_content.get("data_highlights"):
                                data_text = "\n".join([f"📊 {data}" for data in structured_content["data_highlights"]])
                                content_map[placeholder_key] = data_text
                                logger.info(f"       → 分配数据高亮: {len(data_text)}字符")
                            elif structured_content.get("case_studies"):
                                case_text = "\n".join([f"💡 {case}" for case in structured_content["case_studies"]])
                                content_map[placeholder_key] = case_text
                                logger.info(f"       → 分配案例说明: {len(case_text)}字符")
                            else:
                                # 如果没有数据/案例，使用合并的body_content
                                content = self._build_body_content(structured_content, placeholder_type)
                                content_map[placeholder_key] = content
                                logger.info(f"       → 分配正文内容: {len(content)}字符")
                        else:
                            # 其他占位符：使用合并的body_content
                            content = self._build_body_content(structured_content, placeholder_type)
                            content_map[placeholder_key] = content
                            logger.info(f"       → 分配正文内容: {len(content)}字符")
                            logger.info(f"         内容预览: {content[:150]}...")
                    else:
                        # 其他占位符
                        assigned_content = structured_content["main_content"]
                        content_map[placeholder_key] = assigned_content
                        logger.info(f"       → 分配主要内容: {len(assigned_content)}字符")
                        logger.info(f"         内容预览: {assigned_content[:150]}...")
            else:
                # 【修复】对于超出框架PPT范围的幻灯片，使用默认占位符结构
                # 【改进】将不同类型的内容拆分成多个内容块，提升视觉层次
                logger.info(f"     幻灯片{slide_idx}超出框架PPT范围，使用默认占位符结构（拆分内容块）")
                
                # 标题占位符
                title_key = f"slide_{slide_idx}_placeholder_0"
                title_content = structured_content.get("title", "")
                if title_content:
                    content_map[title_key] = title_content
                    logger.info(f"       → 分配标题内容: {title_content[:100]}...")
                
                # 【改进】将内容拆分成多个块，而不是合并成一个
                placeholder_idx = 1
                
                # 1. 关键要点（如果有，单独一个块）
                if structured_content.get("key_points"):
                    key_points_text = "\n".join([f"• {point}" for point in structured_content["key_points"]])
                    key_points_key = f"slide_{slide_idx}_placeholder_{placeholder_idx}"
                    content_map[key_points_key] = key_points_text
                    logger.info(f"       → 分配关键要点块: {len(key_points_text)}字符, {len(structured_content['key_points'])}个要点")
                    placeholder_idx += 1
                
                # 2. 数据高亮（如果有，单独一个块）
                if structured_content.get("data_highlights"):
                    data_text = "\n".join([f"📊 {data}" for data in structured_content["data_highlights"]])
                    data_key = f"slide_{slide_idx}_placeholder_{placeholder_idx}"
                    content_map[data_key] = data_text
                    logger.info(f"       → 分配数据高亮块: {len(data_text)}字符, {len(structured_content['data_highlights'])}个数据")
                    placeholder_idx += 1
                
                # 3. 案例说明（如果有，单独一个块）
                if structured_content.get("case_studies"):
                    case_text = "\n".join([f"💡 {case}" for case in structured_content["case_studies"]])
                    case_key = f"slide_{slide_idx}_placeholder_{placeholder_idx}"
                    content_map[case_key] = case_text
                    logger.info(f"       → 分配案例说明块: {len(case_text)}字符, {len(structured_content['case_studies'])}个案例")
                    placeholder_idx += 1
                
                # 4. 主要内容（如果有且简短，作为概述）
                if structured_content.get("main_content") and len(structured_content["main_content"]) < 100:
                    main_key = f"slide_{slide_idx}_placeholder_{placeholder_idx}"
                    content_map[main_key] = structured_content["main_content"]
                    logger.info(f"       → 分配主要内容块: {len(structured_content['main_content'])}字符")
                    placeholder_idx += 1
                
                # 如果没有以上任何内容，使用合并的body_content作为后备
                if placeholder_idx == 1:  # 只有标题，没有其他内容
                    body_key = f"slide_{slide_idx}_placeholder_1"
                    body_content = self._build_body_content(structured_content, "OBJECT")
                    if body_content:
                        content_map[body_key] = body_content
                        logger.info(f"       → 分配正文内容（后备）: {len(body_content)}字符")
        
        logger.info(f"\n   映射完成，共生成{len(content_map)}个内容映射项")
        logger.info("="*80)
        
        return content_map
    
    def _build_body_content(
        self,
        structured_content: Dict[str, Any],
        placeholder_type: str
    ) -> str:
        """构建正文内容（要点式，不是发言稿）"""
        parts = []
        
        # 1. 关键要点（优先显示，这是PPT的主要内容）
        if structured_content.get("key_points"):
            for point in structured_content["key_points"]:
                parts.append(f"• {point}")
        
        # 2. 数据高亮（如果有，单独列出）
        if structured_content.get("data_highlights"):
            for data in structured_content["data_highlights"]:
                parts.append(f"📊 {data}")
        
        # 3. 案例说明（如果有，简洁列出）
        if structured_content.get("case_studies"):
            for case in structured_content["case_studies"]:
                parts.append(f"💡 {case}")
        
        # 4. 主要内容（如果有且简短，可以作为概述）
        # 但不要放长段落，因为那是发言稿风格
        if structured_content.get("main_content") and len(structured_content["main_content"]) < 100:
            # 只在内容简短时使用
            parts.insert(0, structured_content["main_content"])
        
        return "\n".join(parts) if parts else ""
    
    def _create_structure_from_docx_content(
        self,
        docx_content: str,
        framework_structure: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        从docx内容创建结构数据，用于人类中心化分析
        
        Args:
            docx_content: docx文档的文本内容
            framework_structure: 框架PPT的结构（用于获取幻灯片布局信息）
            
        Returns:
            结构数据字典，格式与EnhancedPPTParser的输出类似
        """
        # 将docx内容按段落分割
        paragraphs = [p.strip() for p in docx_content.split('\n') if p.strip()]
        
        # 【改进】识别标题和章节结构
        # 标题特征：短文本（<50字）、可能包含emoji、数字编号、关键词等
        headings = []
        for i, para in enumerate(paragraphs):
            # 判断是否为标题
            is_heading = False
            heading_level = 0
            
            # 特征1: 短文本（<50字）且可能是标题
            if len(para) < 50:
                # 特征2: 包含emoji（如💎、🚀）
                if any(ord(c) > 127 and c not in '，。、；：！？""''（）【】《》' for c in para[:10]):
                    is_heading = True
                    heading_level = 1
                # 特征3: 数字编号（如"1、"、"2、"、"一、"等）
                elif re.match(r'^[0-9一二三四五六七八九十]+[、.]', para):
                    is_heading = True
                    heading_level = 2
                # 特征4: 包含关键词（如"分析"、"战略"、"路线图"、"回顾"等）
                elif any(kw in para for kw in ['分析', '战略', '路线图', '回顾', '市场', '商业化', '技术', '产品', '文档', '启示', '规划', '能力', '路径']):
                    is_heading = True
                    heading_level = 1
                # 特征5: 主标题（如"技术产品商业化-文档"）
                elif '-' in para or '：' in para or ':' in para:
                    # 如果包含短横线或冒号，且长度较短，可能是标题
                    if len(para) < 30:
                        is_heading = True
                        heading_level = 1
            
            if is_heading:
                headings.append({
                    "index": i,
                    "text": para,
                    "level": heading_level
                })
        
        logger.info(f"--- [PPTFiller]: 识别到{len(headings)}个标题/章节")
        for h in headings[:15]:  # 显示前15个
            logger.info(f"   标题: {h['text'][:50]}")
        
        # 【改进】基于标题将内容分成多个板块，每个板块生成1-2张幻灯片
        slides = []
        framework_slides = framework_structure.get("slides", [])
        framework_placeholders = framework_slides[0].get("placeholders", []) if framework_slides else []
        
        if headings:
            # 基于标题分割内容
            for heading_idx, heading in enumerate(headings):
                start_para_idx = heading["index"]
                end_para_idx = headings[heading_idx + 1]["index"] if heading_idx + 1 < len(headings) else len(paragraphs)
                
                # 获取该板块的段落
                section_paragraphs = paragraphs[start_para_idx:end_para_idx]
                
                # 每个板块生成1-2张幻灯片（根据内容长度）
                # 如果内容超过200字，分成2张幻灯片
                section_text = ' '.join(section_paragraphs)
                num_slides_for_section = 2 if len(section_text) > 200 else 1
                
                paragraphs_per_slide = max(1, len(section_paragraphs) // num_slides_for_section)
                
                for slide_in_section in range(num_slides_for_section):
                    slide_start = slide_in_section * paragraphs_per_slide
                    slide_end = (slide_in_section + 1) * paragraphs_per_slide if slide_in_section < num_slides_for_section - 1 else len(section_paragraphs)
                    slide_paragraphs = section_paragraphs[slide_start:slide_end]
                    
                    # 第一张幻灯片包含标题
                    if slide_in_section == 0:
                        slide_paragraphs = [heading["text"]] + slide_paragraphs
                    
                    # 创建幻灯片结构
                    slide_data = {
                        "slide_index": len(slides),
                        "placeholders": framework_placeholders,
                        "shapes": []
                    }
                    
                    # 为每个段落创建一个shape
                    for para_idx, para_text in enumerate(slide_paragraphs):
                        shape_data = {
                            "shape_id": para_idx,
                            "text": para_text,
                            "shape_type": "paragraph",
                            "is_placeholder": False,
                            "format": {
                                "is_bold": para_idx == 0 and slide_in_section == 0,  # 标题加粗
                                "font_size_pt": 20 if para_idx == 0 and slide_in_section == 0 else 14
                            }
                        }
                        slide_data["shapes"].append(shape_data)
                    
                    slides.append(slide_data)
        else:
            # 如果没有识别到标题，使用原来的策略
            logger.warning("--- [PPTFiller]: 未识别到标题，使用简单分配策略")
            if framework_slides:
                paragraphs_per_slide = max(1, len(paragraphs) // len(framework_slides))
                
                for slide_idx, framework_slide in enumerate(framework_slides):
                    start_idx = slide_idx * paragraphs_per_slide
                    end_idx = (slide_idx + 1) * paragraphs_per_slide if slide_idx < len(framework_slides) - 1 else len(paragraphs)
                    slide_paragraphs = paragraphs[start_idx:end_idx]
                    
                    slide_data = {
                        "slide_index": slide_idx,
                        "placeholders": framework_slide.get("placeholders", []),
                        "shapes": []
                    }
                    
                    for para_idx, para_text in enumerate(slide_paragraphs):
                        shape_data = {
                            "shape_id": para_idx,
                            "text": para_text,
                            "shape_type": "paragraph",
                            "is_placeholder": False
                        }
                        slide_data["shapes"].append(shape_data)
                    
                    slides.append(slide_data)
            else:
                # 如果没有框架，创建一个简单的单幻灯片结构
                slide_data = {
                    "slide_index": 0,
                    "placeholders": [],
                    "shapes": []
                }
                for para_idx, para_text in enumerate(paragraphs):
                    shape_data = {
                        "shape_id": para_idx,
                        "text": para_text,
                        "shape_type": "paragraph",
                        "is_placeholder": False
                    }
                    slide_data["shapes"].append(shape_data)
                slides.append(slide_data)
        
        logger.info(f"--- [PPTFiller]: 基于docx结构生成了{len(slides)}张幻灯片结构（之前只有{len(framework_slides) if framework_slides else 1}张）")
        
        return {
            "slide_count": len(slides),
            "slides": slides
        }
    
    def _generate_output_path(self) -> str:
        """生成输出文件路径"""
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        output_name = f"{self.framework_path.stem}-filled-{timestamp}.pptx"
        output_path = self.framework_path.parent / output_name
        return str(output_path)

