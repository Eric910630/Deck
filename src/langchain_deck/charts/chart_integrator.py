"""
图表整合器
识别可可视化数据，生成图表，并插入到PPT
"""

from typing import List, Dict, Any, Optional
from pathlib import Path
from loguru import logger
from pptx import Presentation
from pptx.util import Cm

from .chart_generator import ChartGenerator
from ..analysis.supporting_materials_analyzer import SupportingMaterialsAnalyzer


class ChartIntegrator:
    """
    图表整合器
    识别可可视化数据，生成图表，并插入到PPT
    """
    
    def __init__(
        self,
        chart_generator: Optional[ChartGenerator] = None,
        materials_analyzer: Optional[SupportingMaterialsAnalyzer] = None
    ):
        """
        初始化图表整合器
        
        Args:
            chart_generator: 图表生成器
            materials_analyzer: 支撑材料分析器
        """
        self.chart_generator = chart_generator or ChartGenerator()
        self.materials_analyzer = materials_analyzer
        logger.info("--- [ChartIntegrator]: 初始化图表整合器")
    
    async def integrate_charts(
        self,
        prs: Presentation,
        data_points: List[Dict[str, Any]],
        output_dir: Optional[Path] = None
    ) -> int:
        """
        整合图表到PPT
        
        Args:
            prs: PPT演示文稿对象
            data_points: 智能识别的数据点列表
            output_dir: 图表输出目录
            
        Returns:
            生成的图表数量
        """
        if not data_points:
            logger.info("--- [ChartIntegrator]: 无数据点，跳过图表生成")
            return 0
        
        logger.info(f"--- [ChartIntegrator]: 开始整合图表，数据点数量: {len(data_points)}")
        
        # 1. 识别可可视化数据
        if self.materials_analyzer:
            chartable_data = await self.materials_analyzer.identify_chartable_data(data_points)
        else:
            # Fallback: 简单判断
            chartable_data = self._simple_identify_chartable_data(data_points)
        
        if not chartable_data:
            logger.info("--- [ChartIntegrator]: 无适合可视化的数据")
            return 0
        
        logger.info(f"--- [ChartIntegrator]: 识别出{len(chartable_data)}个可可视化数据")
        
        # 2. 生成图表并插入到PPT
        chart_count = 0
        for chart_data in chartable_data:
            try:
                chart_path = self._generate_chart(chart_data, output_dir)
                if chart_path and chart_path.exists():
                    self._insert_chart_to_ppt(prs, chart_path, chart_data)
                    chart_count += 1
                    logger.info(f"--- [ChartIntegrator]: 成功插入图表到幻灯片{chart_data.get('slide_index', 0)}")
            except Exception as e:
                logger.warning(f"--- [ChartIntegrator]: 生成/插入图表失败: {e}")
        
        logger.info(f"--- [ChartIntegrator]: 共生成并插入{chart_count}个图表")
        return chart_count
    
    def _generate_chart(
        self,
        chart_data: Dict[str, Any],
        output_dir: Optional[Path] = None
    ) -> Optional[Path]:
        """
        生成图表
        
        Args:
            chart_data: 图表数据
            output_dir: 输出目录
            
        Returns:
            图表文件路径
        """
        chart_type = chart_data.get("chart_type", "bar")
        data = chart_data.get("data", [])
        title = chart_data.get("title", "图表")
        
        if not data:
            return None
        
        try:
            # 转换为图表生成器需要的格式
            chart_config = {
                "type": chart_type,
                "title": title,
                "data": data,
                "x_axis": chart_data.get("x_axis", ""),
                "y_axis": chart_data.get("y_axis", "")
            }
            
            # 生成图表（使用现有的方法）
            # 确保数据格式正确
            chart_data_list = []
            for item in data:
                if isinstance(item, dict):
                    chart_data_list.append(item)
                else:
                    # 如果是简单格式，转换为字典
                    chart_data_list.append({"label": str(item.get("label", "")), "value": float(item.get("value", 0))})
            
            if chart_type == "bar":
                chart_path = self.chart_generator.generate_bar_chart(
                    data=chart_data_list,
                    x_key="label",
                    y_key="value",
                    title=title
                )
            elif chart_type == "line":
                chart_path = self.chart_generator.generate_line_chart(
                    data=chart_data_list,
                    x_key="label",
                    y_key="value",
                    title=title
                )
            elif chart_type == "pie":
                chart_path = self.chart_generator.generate_pie_chart(
                    data=chart_data_list,
                    label_key="label",
                    value_key="value",
                    title=title
                )
            else:
                # 默认使用柱状图
                chart_path = self.chart_generator.generate_bar_chart(
                    data=chart_data_list,
                    x_key="label",
                    y_key="value",
                    title=title
                )
            
            if chart_path:
                chart_path_obj = Path(chart_path) if isinstance(chart_path, str) else chart_path
                return chart_path_obj
            return None
            
        except Exception as e:
            logger.error(f"--- [ChartIntegrator]: 生成图表失败: {e}")
            import traceback
            logger.debug(traceback.format_exc())
            return None
    
    def _insert_chart_to_ppt(
        self,
        prs: Presentation,
        chart_path: Path,
        chart_data: Dict[str, Any]
    ):
        """
        将图表插入到PPT
        
        Args:
            prs: PPT演示文稿对象
            chart_path: 图表文件路径
            chart_data: 图表数据（包含位置信息）
        """
        slide_index = chart_data.get("slide_index", 0)
        position = chart_data.get("recommended_position", {})
        
        # 确保幻灯片存在
        if slide_index >= len(prs.slides):
            logger.warning(f"--- [ChartIntegrator]: 幻灯片{slide_index}不存在，跳过")
            return
        
        slide = prs.slides[slide_index]
        
        # 获取位置信息（默认值）
        x = Cm(position.get("x", 10))
        y = Cm(position.get("y", 5))
        width = Cm(position.get("width", 15))
        height = Cm(position.get("height", 8))
        
        # 插入图表图片
        slide.shapes.add_picture(
            str(chart_path),
            x,
            y,
            width,
            height
        )
        
        logger.info(f"--- [ChartIntegrator]: 图表插入到幻灯片{slide_index}，位置: ({x}, {y}), 尺寸: {width} × {height}")
    
    def _simple_identify_chartable_data(
        self,
        data_points: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        """
        简单识别可可视化数据（Fallback方法）
        
        Args:
            data_points: 数据点列表
            
        Returns:
            可可视化数据列表
        """
        # 简单策略：如果有3个或以上数据点，生成柱状图
        if len(data_points) >= 3:
            chart_data = {
                "chart_type": "bar",
                "data": [
                    {
                        "label": dp.get("label", dp.get("value", "")),
                        "value": self._extract_numeric_value(dp.get("value", ""))
                    }
                    for dp in data_points[:5]  # 最多5个数据点
                ],
                "title": "数据对比",
                "x_axis": "指标",
                "y_axis": "数值",
                "slide_index": data_points[0].get("slide_index", 0),
                "recommended_position": {
                    "x": 10,
                    "y": 5,
                    "width": 15,
                    "height": 8
                }
            }
            return [chart_data]
        
        return []
    
    def _extract_numeric_value(self, value_str: str) -> float:
        """
        从字符串中提取数值
        
        Args:
            value_str: 数值字符串（如"40-60%"、"50%"等）
            
        Returns:
            提取的数值
        """
        import re
        
        # 提取数字
        numbers = re.findall(r'\d+\.?\d*', value_str)
        if numbers:
            # 如果有范围，取平均值
            if len(numbers) >= 2:
                return (float(numbers[0]) + float(numbers[1])) / 2
            else:
                return float(numbers[0])
        
        return 0.0

