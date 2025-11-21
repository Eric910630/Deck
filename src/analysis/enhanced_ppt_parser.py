"""
增强版PPT解析器
添加格式提取、段落分析、列表识别等功能
"""

from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from loguru import logger
import re


class EnhancedPPTParser:
    """
    增强版PPT解析器
    在基础解析的基础上，增加格式提取、段落分析、列表识别等功能
    """
    
    def __init__(self, ppt_path: str):
        """初始化增强解析器"""
        self.ppt_path = Path(ppt_path)
        if not self.ppt_path.exists():
            raise FileNotFoundError(f"PPT file not found: {ppt_path}")
        
        self.prs = Presentation(str(self.ppt_path))
        logger.info(f"--- [EnhancedPPTParser]: Loaded PPT: {self.ppt_path}")
    
    def _extract_format_info(self, shape) -> Dict[str, Any]:
        """
        提取形状的格式信息
        
        Args:
            shape: PPT形状对象
            
        Returns:
            格式信息字典
        """
        format_info = {
            "font_name": None,
            "font_size_pt": None,
            "font_color": None,
            "is_bold": False,
            "is_italic": False,
            "is_underline": False,
            "alignment": None,
            "line_spacing": None,
            "left_indent_pt": None,
            "first_line_indent_pt": None
        }
        
        if hasattr(shape, "text_frame"):
            # 收集所有运行（runs）的格式信息
            font_sizes = []
            font_names = []
            font_colors = []
            has_bold = False
            has_italic = False
            has_underline = False
            
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        font_names.append(run.font.name)
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)
                    try:
                        if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                            font_colors.append(str(run.font.color.rgb))
                    except:
                        pass
                    if hasattr(run, 'bold') and run.bold:
                        has_bold = True
                    if hasattr(run, 'italic') and run.italic:
                        has_italic = True
                    if hasattr(run, 'underline') and run.underline:
                        has_underline = True
                
                # 段落格式
                try:
                    pf = para.paragraph_format
                    if para.alignment:
                        format_info["alignment"] = str(para.alignment)
                    if pf.line_spacing:
                        format_info["line_spacing"] = pf.line_spacing
                    if pf.left_indent:
                        format_info["left_indent_pt"] = pf.left_indent.pt
                    if pf.first_line_indent:
                        format_info["first_line_indent_pt"] = pf.first_line_indent.pt
                except:
                    pass
            
            # 取最常见的值
            if font_names:
                format_info["font_name"] = max(set(font_names), key=font_names.count)
            if font_sizes:
                format_info["font_size_pt"] = max(font_sizes)  # 取最大字号
            if font_colors:
                format_info["font_color"] = max(set(font_colors), key=font_colors.count)
            
            format_info["is_bold"] = has_bold
            format_info["is_italic"] = has_italic
            format_info["is_underline"] = has_underline
        
        return format_info
    
    def extract_structure_enhanced(self) -> Dict[str, Any]:
        """
        提取增强的结构信息（包含格式信息）
        
        Returns:
            增强的结构信息字典
        """
        structure = {
            "slides": [],
            "slide_count": len(self.prs.slides),
            "slide_width": float(self.prs.slide_width) / 360000,
            "slide_height": float(self.prs.slide_height) / 360000
        }
        
        for idx, slide in enumerate(self.prs.slides):
            slide_info = {
                "slide_index": idx,
                "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown",
                "shapes": [],
                "placeholders": [],
                "format_statistics": {
                    "font_sizes": set(),
                    "font_names": set(),
                    "bold_count": 0,
                    "total_shapes": 0
                }
            }
            
            for shape in slide.shapes:
                shape_info = self._extract_shape_info_enhanced(shape, idx)
                if shape_info:
                    slide_info["shapes"].append(shape_info)
                    
                    # 统计格式信息
                    format_info = shape_info.get("format", {})
                    if format_info.get("font_size_pt"):
                        slide_info["format_statistics"]["font_sizes"].add(format_info["font_size_pt"])
                    if format_info.get("font_name"):
                        slide_info["format_statistics"]["font_names"].add(format_info["font_name"])
                    if format_info.get("is_bold"):
                        slide_info["format_statistics"]["bold_count"] += 1
                    
                    slide_info["format_statistics"]["total_shapes"] += 1
                    
                    if shape.is_placeholder:
                        slide_info["placeholders"].append(shape_info)
            
            # 转换set为list
            slide_info["format_statistics"]["font_sizes"] = sorted(list(slide_info["format_statistics"]["font_sizes"]))
            slide_info["format_statistics"]["font_names"] = list(slide_info["format_statistics"]["font_names"])
            
            structure["slides"].append(slide_info)
        
        logger.info(f"--- [EnhancedPPTParser]: Extracted enhanced structure from {len(structure['slides'])} slides")
        return structure
    
    def _extract_shape_info_enhanced(self, shape, slide_index: int) -> Optional[Dict[str, Any]]:
        """提取增强的形状信息（包含格式）"""
        try:
            shape_info = {
                "type": self._get_shape_type(shape),
                "shape_id": shape.shape_id,
                "left": float(shape.left) / 360000,
                "top": float(shape.top) / 360000,
                "width": float(shape.width) / 360000,
                "height": float(shape.height) / 360000,
                "is_placeholder": shape.is_placeholder,
                "format": self._extract_format_info(shape)  # 新增：格式信息
            }
            
            if shape.is_placeholder:
                try:
                    shape_info["placeholder_id"] = shape.placeholder_format.idx
                    shape_info["placeholder_type"] = str(shape.placeholder_format.type)
                except:
                    pass
            
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text.strip()
                if text:
                    shape_info["text"] = text
                    shape_info["has_text"] = True
                else:
                    shape_info["has_text"] = False
                    shape_info["text"] = ""
            else:
                shape_info["has_text"] = False
                shape_info["text"] = ""
            
            return shape_info
            
        except Exception as e:
            logger.warning(f"--- [EnhancedPPTParser]: Failed to extract shape info: {e}")
            return None
    
    def _get_shape_type(self, shape) -> str:
        """获取形状类型名称"""
        try:
            shape_type = shape.shape_type
            type_names = {
                MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
                MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
                MSO_SHAPE_TYPE.PICTURE: "picture",
                MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
                MSO_SHAPE_TYPE.GROUP: "group",
                MSO_SHAPE_TYPE.TABLE: "table",
                MSO_SHAPE_TYPE.MEDIA: "media"
            }
            return type_names.get(shape_type, "unknown")
        except:
            return "unknown"
    
    def extract_paragraph_structure(self) -> List[Dict[str, Any]]:
        """
        提取段落级别的结构
        
        Returns:
            段落结构列表
        """
        paragraph_structure = []
        
        for slide_idx, slide in enumerate(self.prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame"):
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        para_info = {
                            "slide_index": slide_idx,
                            "shape_index": shape_idx,
                            "paragraph_index": para_idx,
                            "text": para.text.strip(),
                            "format": self._extract_paragraph_format(para),
                            "runs": []
                        }
                        
                        # 分析文本运行
                        for run in para.runs:
                            run_info = {
                                "text": run.text,
                                "format": self._extract_run_format(run)
                            }
                            para_info["runs"].append(run_info)
                        
                        paragraph_structure.append(para_info)
        
        logger.info(f"--- [EnhancedPPTParser]: Extracted {len(paragraph_structure)} paragraphs")
        return paragraph_structure
    
    def _extract_paragraph_format(self, para) -> Dict[str, Any]:
        """提取段落格式"""
        try:
            pf = para.paragraph_format
            return {
                "alignment": str(para.alignment) if para.alignment else None,
                "left_indent_pt": pf.left_indent.pt if pf.left_indent else 0,
                "first_line_indent_pt": pf.first_line_indent.pt if pf.first_line_indent else 0,
                "space_before_pt": pf.space_before.pt if pf.space_before else 0,
                "space_after_pt": pf.space_after.pt if pf.space_after else 0,
                "line_spacing": str(pf.line_spacing) if pf.line_spacing else None
            }
        except Exception as e:
            logger.warning(f"--- [EnhancedPPTParser]: Failed to extract paragraph format: {e}")
            return {
                "alignment": None,
                "left_indent_pt": 0,
                "first_line_indent_pt": 0,
                "space_before_pt": 0,
                "space_after_pt": 0,
                "line_spacing": None
            }
    
    def _extract_run_format(self, run) -> Dict[str, Any]:
        """提取文本运行格式"""
        font_color = None
        try:
            if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                font_color = str(run.font.color.rgb)
        except:
            pass
        
        return {
            "font_name": run.font.name if run.font.name else None,
            "font_size_pt": run.font.size.pt if run.font.size else None,
            "font_color": font_color,
            "is_bold": run.bold if hasattr(run, 'bold') and run.bold is not None else False,
            "is_italic": run.italic if hasattr(run, 'italic') and run.italic is not None else False,
            "is_underline": run.underline if hasattr(run, 'underline') and run.underline is not None else False
        }
    
    def extract_list_structure(self) -> Dict[str, Any]:
        """
        提取列表结构
        
        Returns:
            列表结构字典
        """
        lists = {
            "numbered_lists": [],
            "bullet_lists": [],
            "indented_items": []
        }
        
        for slide_idx, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        
                        # 检查编号列表
                        numbered_match = re.match(r'^[\d一二三四五六七八九十]+[\.、]\s*(.+)', text)
                        if numbered_match:
                            lists["numbered_lists"].append({
                                "slide_index": slide_idx,
                                "text": text,
                                "number": text[0] if text else "",
                                "content": numbered_match.group(1) if len(numbered_match.groups()) > 0 else text,
                                "format": self._extract_paragraph_format(para)
                            })
                        
                        # 检查项目符号
                        bullet_match = re.match(r'^[•·▪▫○●■□]\s*(.+)', text)
                        if bullet_match:
                            lists["bullet_lists"].append({
                                "slide_index": slide_idx,
                                "text": text,
                                "bullet": text[0],
                                "content": bullet_match.group(1),
                                "format": self._extract_paragraph_format(para)
                            })
                        
                        # 检查缩进
                        try:
                            pf = para.paragraph_format
                            if pf.left_indent and pf.left_indent.pt > 0:
                                lists["indented_items"].append({
                                    "slide_index": slide_idx,
                                    "text": text[:100],
                                    "indent_pt": pf.left_indent.pt,
                                    "format": self._extract_paragraph_format(para)
                                })
                        except:
                            pass
        
        logger.info(f"--- [EnhancedPPTParser]: Extracted {len(lists['numbered_lists'])} numbered lists, {len(lists['bullet_lists'])} bullet lists, {len(lists['indented_items'])} indented items")
        return lists
    
    def extract_table_structure(self) -> List[Dict[str, Any]]:
        """
        提取表格结构
        
        Returns:
            表格结构列表
        """
        tables = []
        
        for slide_idx, slide in enumerate(self.prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_info = {
                        "slide_index": slide_idx,
                        "rows": len(shape.table.rows),
                        "columns": len(shape.table.columns) if shape.table.rows else 0,
                        "cells": [],
                        "structure": []
                    }
                    
                    for row_idx, row in enumerate(shape.table.rows):
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            row_data.append(cell_text)
                            table_info["cells"].append({
                                "row": row_idx,
                                "column": row.cells.index(cell),
                                "text": cell_text
                            })
                        table_info["structure"].append(row_data)
                    
                    # 识别表头
                    if table_info["structure"]:
                        table_info["header_row"] = table_info["structure"][0]
                        table_info["data_rows"] = table_info["structure"][1:]
                    
                    tables.append(table_info)
        
        logger.info(f"--- [EnhancedPPTParser]: Extracted {len(tables)} tables")
        return tables

