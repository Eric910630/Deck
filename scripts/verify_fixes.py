#!/usr/bin/env python3
"""
验证修复效果的脚本
检查生成的PPT是否符合16:9和Ant Design规范
"""

from pptx import Presentation
from pptx.util import Pt
from pathlib import Path
import glob

def verify_ppt(ppt_path: str):
    """验证PPT文件"""
    print(f"\n{'='*80}")
    print(f"验证文件: {ppt_path}")
    print('='*80)
    
    prs = Presentation(ppt_path)
    
    # 1. 检查尺寸
    print("\n【1. 尺寸检查】")
    width_emu = prs.slide_width
    height_emu = prs.slide_height
    width_cm = float(width_emu) / 360000
    height_cm = float(height_emu) / 360000
    ratio = width_cm / height_cm
    is_16_9 = abs(ratio - 16/9) < 0.1
    
    print(f"  宽度: {width_cm:.2f}cm ({width_emu:,} EMU)")
    print(f"  高度: {height_cm:.2f}cm ({height_emu:,} EMU)")
    print(f"  宽高比: {ratio:.2f} (16:9={16/9:.2f})")
    print(f"  {'✅' if is_16_9 else '❌'} 是否为16:9: {is_16_9}")
    
    # 2. 检查设计规范
    print("\n【2. 设计规范检查】")
    if prs.slides:
        slide = prs.slides[0]
        placeholder_count = 0
        
        for shape in slide.shapes:
            if shape.is_placeholder and hasattr(shape, 'text_frame'):
                placeholder_count += 1
                print(f"\n  占位符 {placeholder_count}:")
                
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs:
                        font = para.runs[0].font
                        
                        # 字体
                        font_name = font.name or "未设置"
                        print(f"    字体: {font_name}")
                        
                        # 字号（EMU转pt）
                        if font.size:
                            size_pt = float(font.size) / 12700
                            print(f"    字号: {size_pt:.0f}pt (EMU: {font.size:,})")
                            
                            # 验证字号是否符合规范
                            if placeholder_count == 1:
                                # 第一个占位符应该是标题（38pt）
                                is_correct = abs(size_pt - 38) < 1
                                print(f"    {'✅' if is_correct else '❌'} 标题字号检查: {is_correct} (期望38pt)")
                            else:
                                # 其他占位符应该是正文（14pt）
                                is_correct = abs(size_pt - 14) < 1
                                print(f"    {'✅' if is_correct else '❌'} 正文字号检查: {is_correct} (期望14pt)")
                        else:
                            print(f"    ❌ 字号: 未设置")
                        
                        # 加粗
                        print(f"    加粗: {font.bold}")
                        
                        # 颜色
                        if font.color and font.color.rgb:
                            rgb = font.color.rgb
                            # RGBColor对象可以转换为整数
                            try:
                                if hasattr(rgb, '__int__'):
                                    rgb_int = int(rgb)
                                    r = (rgb_int >> 16) & 0xFF
                                    g = (rgb_int >> 8) & 0xFF
                                    b = rgb_int & 0xFF
                                elif isinstance(rgb, int):
                                    r = (rgb >> 16) & 0xFF
                                    g = (rgb >> 8) & 0xFF
                                    b = rgb & 0xFF
                                else:
                                    # 尝试从字符串解析（格式可能是"262626"）
                                    rgb_str = str(rgb)
                                    if len(rgb_str) == 6 and rgb_str.isdigit():
                                        # 直接是hex字符串
                                        r = int(rgb_str[0:2], 16)
                                        g = int(rgb_str[2:4], 16)
                                        b = int(rgb_str[4:6], 16)
                                    else:
                                        r, g, b = 0, 0, 0
                            except Exception as e:
                                print(f"    颜色解析错误: {e}")
                                r, g, b = 0, 0, 0
                            
                            hex_color = f"#{r:02x}{g:02x}{b:02x}"
                            print(f"    颜色: RGB({r}, {g}, {b}) = {hex_color}")
                            
                            # 验证颜色是否符合规范（#262626）
                            expected_rgb = (38, 38, 38)
                            is_correct = (r, g, b) == expected_rgb
                            print(f"    {'✅' if is_correct else '❌'} 颜色检查: {is_correct} (期望#262626)")
                        else:
                            print(f"    ❌ 颜色: 未设置")
    
    # 3. 总结
    print("\n【3. 验证总结】")
    all_checks = [
        ("16:9尺寸", is_16_9),
    ]
    
    for check_name, check_result in all_checks:
        status = "✅ 通过" if check_result else "❌ 失败"
        print(f"  {check_name}: {status}")
    
    return is_16_9


if __name__ == "__main__":
    # 查找最新生成的PPT文件
    files = sorted(glob.glob('demo_filled-filled-*.pptx'))
    
    if not files:
        print("未找到生成的PPT文件")
        print("请先运行: python test_demo_framework.py")
    else:
        print(f"找到 {len(files)} 个生成的PPT文件")
        print("验证最新的文件...")
        
        latest = files[-1]
        verify_ppt(latest)
        
        print("\n" + "="*80)
        print("验证完成！")
        print("="*80)

