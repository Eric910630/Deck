#!/usr/bin/env python3
"""
创建示例PPT框架文件
用于演示框架填充功能
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_framework_ppt():
    """创建一个包含占位符的PPT框架（16:9横版）"""
    from pptx.util import Cm
    
    prs = Presentation()
    # 16:9 横版尺寸 (33.867cm x 19.05cm)
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    
    # 第1张幻灯片 - 标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    
    # 获取标题和副标题占位符
    title_shape = slide1.shapes.title
    subtitle_shape = slide1.placeholders[1]
    
    # 设置占位符提示文本（这些会被LLM生成的内容替换）
    title_shape.text = "[标题占位符]"
    subtitle_shape.text = "[副标题占位符]"
    
    # 第2张幻灯片 - 内容页
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    
    title_shape2 = slide2.shapes.title
    content_shape2 = slide2.placeholders[1]
    
    title_shape2.text = "[内容页标题]"
    content_shape2.text = "[内容占位符]\n\n这里可以添加详细内容..."
    
    # 第3张幻灯片 - 空白页（用于自定义布局）
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # 添加标题文本框
    title_box = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "[自定义标题]"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # 添加内容文本框
    content_box = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    content_frame = content_box.text_frame
    content_frame.text = "[自定义内容占位符]\n\n可以在这里添加更多内容..."
    content_frame.paragraphs[0].font.size = Pt(18)
    
    # 保存文件
    output_path = "framework_template.pptx"
    prs.save(output_path)
    print(f"✓ 框架PPT已创建: {output_path}")
    print(f"  - 包含 {len(prs.slides)} 张幻灯片")
    print(f"  - 包含多个占位符，等待填充")
    
    return output_path

if __name__ == "__main__":
    create_framework_ppt()

